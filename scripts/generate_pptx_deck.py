"""
Factory-Vision MES - Perfected Executive PowerPoint Presentation (python-pptx)
Strictly adheres to Docs/Factory-Vision-MES-Product-Deck-Guideline.md
Style: McKinsey-inspired consulting presentation + modern B2B SaaS product
Theme: Pristine Light Mode with high-resolution screenshot evidence
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SCREENSHOTS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "Docs", "screenshots"))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Perfected Color Palette
COLOR_NAVY = RGBColor(0x12, 0x26, 0x3F)        # #12263F - Deep Executive Navy
COLOR_BLUE = RGBColor(0x1E, 0x5C, 0xA0)        # #1E5CA0 - Corporate Cobalt Blue
COLOR_TEAL = RGBColor(0x0E, 0x74, 0x90)        # #0E7490 - Execution Cyan/Teal
COLOR_GREEN = RGBColor(0x15, 0x80, 0x3D)       # #15803D - Success / Output Green
COLOR_AMBER = RGBColor(0xB4, 0x53, 0x09)       # #B45309 - Warning Amber
COLOR_RED = RGBColor(0xBE, 0x12, 0x3C)         # #BE123C - Downtime & Reject Red
COLOR_BG_CARD = RGBColor(0xF8, 0xFA, 0xFC)     # #F8FAFC - Light Slate Card Background
COLOR_BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # #FFFFFF - Crisp White
COLOR_BORDER = RGBColor(0xCB, 0xD5, 0xE1)      # #CBD5E1 - Slate Border
COLOR_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B - High Contrast Text
COLOR_MED_TEXT = RGBColor(0x47, 0x55, 0x69)    # #475569 - Secondary Body Text
COLOR_MUTED_TEXT = RGBColor(0x64, 0x74, 0x8B)  # #64748B - Labels & Subtitles

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def set_para_font(p, name=FONT_BODY, size_pt=10, color=COLOR_DARK_TEXT, bold=False):
    p.font.name = name
    p.font.size = Pt(size_pt)
    p.font.color.rgb = color
    p.font.bold = bold

def add_header(slide, section_label, title_text, subtitle_text=None):
    """Standardized McKinsey slide header with category kicker, bold title, and takeaway subtitle."""
    # Kicker
    tx_kicker = slide.shapes.add_textbox(Inches(0.65), Inches(0.38), Inches(12.033), Inches(0.25))
    tf_k = tx_kicker.text_frame
    tf_k.word_wrap = True
    tf_k.margin_left = tf_k.margin_top = tf_k.margin_right = tf_k.margin_bottom = 0
    pk = tf_k.paragraphs[0]
    pk.text = f"■  {section_label.upper()}"
    set_para_font(pk, FONT_HEADING, 8.5, COLOR_BLUE, bold=True)

    # Title
    tx_title = slide.shapes.add_textbox(Inches(0.65), Inches(0.65), Inches(12.033), Inches(0.42))
    tf_t = tx_title.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    pt = tf_t.paragraphs[0]
    pt.text = title_text
    set_para_font(pt, FONT_HEADING, 17.5, COLOR_NAVY, bold=True)

    # Subtitle
    if subtitle_text:
        tx_sub = slide.shapes.add_textbox(Inches(0.65), Inches(1.10), Inches(12.033), Inches(0.32))
        tf_s = tx_sub.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
        ps = tf_s.paragraphs[0]
        ps.text = subtitle_text
        set_para_font(ps, FONT_BODY, 10.5, COLOR_MED_TEXT, bold=False)

def add_footer(slide, page_num):
    """Standardized bottom brand and page numbering."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.05), Inches(12.033), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BORDER
    line.line.fill.background()

    tx_brand = slide.shapes.add_textbox(Inches(0.65), Inches(7.12), Inches(6.0), Inches(0.25))
    tf_b = tx_brand.text_frame
    tf_b.margin_left = tf_b.margin_top = 0
    pb = tf_b.paragraphs[0]
    pb.text = "FACTORY-VISION  |  Manufacturing Execution System Product Deck"
    set_para_font(pb, FONT_BODY, 8.5, COLOR_MUTED_TEXT)

    tx_page = slide.shapes.add_textbox(Inches(11.5), Inches(7.12), Inches(1.18), Inches(0.25))
    tf_p = tx_page.text_frame
    tf_p.margin_left = tf_p.margin_top = 0
    pp = tf_p.paragraphs[0]
    pp.text = f"{page_num:02d} / 22"
    pp.alignment = PP_ALIGN.RIGHT
    set_para_font(pp, FONT_BODY, 8.5, COLOR_MUTED_TEXT, bold=True)

def create_card(slide, left, top, width, height, bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

# ==========================================
# SLIDE 01: COVER
# ==========================================
def make_slide_01():
    slide = prs.slides.add_slide(blank_layout)
    
    # Top Accent Ribbon
    top_ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_ribbon.fill.solid()
    top_ribbon.fill.fore_color.rgb = COLOR_BLUE
    top_ribbon.line.fill.background()

    # Kicker Badge
    k_badge = create_card(slide, Inches(0.8), Inches(1.2), Inches(2.2), Inches(0.35), COLOR_BG_CARD, COLOR_BORDER)
    tb_k = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(2.2), Inches(0.35))
    tf_k = tb_k.text_frame
    tf_k.vertical_anchor = MSO_ANCHOR.MIDDLE
    pk = tf_k.paragraphs[0]
    pk.text = "PRODUCT DECK 2026"
    pk.alignment = PP_ALIGN.CENTER
    set_para_font(pk, FONT_HEADING, 8.5, COLOR_BLUE, bold=True)

    # Main Title
    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.5), Inches(1.7))
    tf_t = tx_title.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = 0
    p1 = tf_t.paragraphs[0]
    p1.text = "Factory-Vision"
    set_para_font(p1, FONT_HEADING, 42, COLOR_NAVY, bold=True)

    p2 = tf_t.add_paragraph()
    p2.text = "Manufacturing Execution System"
    set_para_font(p2, FONT_HEADING, 26, COLOR_BLUE, bold=True)

    # Subtitle
    tx_sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10.5), Inches(0.8))
    tf_s = tx_sub.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = 0
    ps = tf_s.paragraphs[0]
    ps.text = "Menghubungkan production planning, shop floor execution, dan production performance dalam satu platform digital terpadu untuk pabrik modern Indonesia."
    set_para_font(ps, FONT_BODY, 13.5, COLOR_MED_TEXT)

    # 3 Pillars
    pillars = [
        ("01. PLAN", "Production Orders, Work Orders, Scheduling, Routing & Lot Allocation", COLOR_BLUE),
        ("02. EXECUTE", "Operator Touch Terminal, Telemetri Mesin, Work In Progress & Shift Control", COLOR_TEAL),
        ("03. MEASURE", "OEE Real-time, Downtime Pareto Analytics, Target vs Actual & BI Cockpit", COLOR_GREEN)
    ]
    card_w = Inches(3.75)
    card_h = Inches(1.85)
    for i, (title, desc, accent) in enumerate(pillars):
        x = Inches(0.8 + i * 4.05)
        y = Inches(4.7)
        card = create_card(slide, x, y, card_w, card_h, COLOR_BG_CARD, COLOR_BORDER)
        
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.3))
        tfr = tb.text_frame
        tfr.word_wrap = True
        tfr.margin_left = tfr.margin_top = 0
        pt = tfr.paragraphs[0]
        pt.text = title
        set_para_font(pt, FONT_HEADING, 13.5, COLOR_NAVY, bold=True)
        
        pd = tfr.add_paragraph()
        pd.text = desc
        pd.space_before = Pt(4)
        set_para_font(pd, FONT_BODY, 10, COLOR_MED_TEXT)

    add_footer(slide, 1)

# ==========================================
# SLIDE 02: CHALLENGES
# ==========================================
def make_slide_02():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "01. BUSINESS CONTEXT",
        "Apa yang Sering Menjadi Tantangan di Operasional Produksi?",
        "Sebagian besar pabrik menghadapi kendala yang berakar dari data yang tersebar dan keterlambatan umpan balik lantai produksi."
    )

    problems = [
        ("Data Produksi Tersebar", "Data tersebar di catatan kertas, formulir fisik, spreadsheet Excel, dan grup chat, menyulitkan konsolidasi data operasional."),
        ("Actual Output Terlambat", "Kondisi output dan scrap baru diketahui di akhir shift atau keesokan harinya saat rekap admin selesai dikerjakan."),
        ("Pencatatan Manual Rentan Error", "Operator menghabiskan banyak waktu mencatat hasil kerja di kertas yang rawan salah tulis, hilang, atau tidak terbaca."),
        ("Akar Masalah Downtime Tidak Jelas", "Mesin mati tanpa pencatatan reason code dan durasi presisi, sehingga tim maintenance sulit menentukan prioritas perbaikan aset."),
        ("Target vs Actual Sulit Dipantau", "Manajemen kesulitan membandingkan gap produksi secara instan untuk mendeteksi order yang berisiko terlambat dikirim."),
        ("Pengambilan Keputusan Lambat", "Manajemen membutuhkan waktu lama untuk menyusun tindakan korektif karena ketiadaan visibilitas operasional real-time.")
    ]

    card_w = Inches(3.8)
    card_h = Inches(2.28)
    
    for i, (title, desc) in enumerate(problems):
        row = i // 3
        col = i % 3
        x = Inches(0.65 + col * 4.1)
        y = Inches(1.70 + row * 2.52)

        card = create_card(slide, x, y, card_w, card_h, COLOR_BG_WHITE, COLOR_BORDER)
        
        # Red top accent strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.04))
        strip.fill.solid()
        strip.fill.fore_color.rgb = COLOR_RED
        strip.line.fill.background()

        # Number Badge
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.18), y + Inches(0.16), Inches(0.48), Inches(0.28))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_BG_CARD
        num_box.line.color.rgb = COLOR_BORDER
        p_num = num_box.text_frame.paragraphs[0]
        p_num.text = f"0{i+1}"
        p_num.alignment = PP_ALIGN.CENTER
        set_para_font(p_num, FONT_HEADING, 9.5, COLOR_RED, bold=True)

        tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.55), card_w - Inches(0.36), card_h - Inches(0.65))
        tfr = tb.text_frame
        tfr.word_wrap = True
        tfr.margin_left = tfr.margin_top = 0
        
        pt = tfr.paragraphs[0]
        pt.text = title
        set_para_font(pt, FONT_HEADING, 11.5, COLOR_NAVY, bold=True)
        
        pd = tfr.add_paragraph()
        pd.text = desc
        pd.space_before = Pt(4)
        set_para_font(pd, FONT_BODY, 9.5, COLOR_MED_TEXT)

    add_footer(slide, 2)

# ==========================================
# SLIDE 03: HOW MES CONNECTS PRODUCTION
# ==========================================
def make_slide_03():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "02. MES OVERVIEW",
        "Factory-Vision Menghubungkan Planning dengan Kondisi Aktual",
        "Aliran data dua arah yang mulus memastikan setiap rencana kerja tereksekusi dan terukur secara presisi di lantai produksi."
    )

    steps = [
        ("01", "Planning", "PPIC merancang target & jadwal produksi pesanan."),
        ("02", "Prod Order", "Pesanan diubah menjadi batch order resmi."),
        ("03", "Work Order", "Tugas dipecah per stasiun & mesin kerja."),
        ("04", "Shop Floor", "Operator menjalankan via terminal digital sentuh."),
        ("05", "Prod Data", "Good, reject, & downtime tercatat real-time."),
        ("06", "OEE & BI", "Sistem mengagregasi availability, speed, & quality."),
        ("07", "Decision", "Manajemen mengambil tindakan mitigasi berbasis data.")
    ]

    card_w = Inches(1.58)
    card_h = Inches(4.35)
    
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.65 + i * 1.74)
        y = Inches(1.70)

        card = create_card(slide, x, y, card_w, card_h, COLOR_BG_WHITE, COLOR_BORDER)
        
        step_color = COLOR_NAVY if i in [0, 6] else (COLOR_BLUE if i in [1, 2] else (COLOR_TEAL if i == 3 else COLOR_GREEN))
        step_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.48))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = step_color
        step_box.line.fill.background()
        p_step = step_box.text_frame.paragraphs[0]
        p_step.text = f"STEP {num}"
        p_step.alignment = PP_ALIGN.CENTER
        set_para_font(p_step, FONT_HEADING, 9, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

        tb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.65), card_w - Inches(0.24), Inches(3.4))
        tfr = tb.text_frame
        tfr.word_wrap = True
        tfr.margin_left = tfr.margin_top = 0
        
        pt = tfr.paragraphs[0]
        pt.text = title
        set_para_font(pt, FONT_HEADING, 11, COLOR_NAVY, bold=True)

        pd = tfr.add_paragraph()
        pd.text = desc
        pd.space_before = Pt(6)
        set_para_font(pd, FONT_BODY, 9, COLOR_MED_TEXT)

    summary_bar = create_card(slide, Inches(0.65), Inches(6.20), Inches(12.033), Inches(0.68), COLOR_BG_CARD, COLOR_BLUE)
    tb_sum = slide.shapes.add_textbox(Inches(0.85), Inches(6.28), Inches(11.6), Inches(0.5))
    tf_sum = tb_sum.text_frame
    tf_sum.word_wrap = True
    tf_sum.margin_left = tf_sum.margin_top = 0
    p_sum = tf_sum.paragraphs[0]
    p_sum.text = "Keunggulan Utama: Menghilangkan jeda waktu informasi antara kantor perencanaan dan stasiun kerja operator secara instan tanpa rekap berganda."
    set_para_font(p_sum, FONT_BODY, 10, COLOR_NAVY, bold=True)

    add_footer(slide, 3)

# ==========================================
# SLIDE 04: 4 PILLARS
# ==========================================
def make_slide_04():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "02. MES OVERVIEW",
        "Satu Platform untuk Menjalankan, Menangkap, dan Memahami Produksi",
        "Arsitektur 4 pilar terpadu yang mencakup seluruh siklus hidup manufaktur dari hulu ke hilir."
    )

    pillars = [
        ("01. PLAN", "Perencanaan & Kontrol", [
            "Production Order Management",
            "Work Order Decomposition",
            "Production Scheduling & Routing",
            "Machine & Line Allocation",
            "Batch & Lot Management"
        ], COLOR_BLUE),
        ("02. EXECUTE", "Eksekusi Shop Floor", [
            "Operator Terminal Touch UI",
            "Machine Execution Dispatch",
            "Start, Pause, Resume, Stop Controls",
            "Shift Workflow Management",
            "Real-time Work In Progress (WIP)"
        ], COLOR_TEAL),
        ("03. CAPTURE", "Penangkapan Data Akurat", [
            "Good Quantity Counting",
            "Reject & Scrap Recording",
            "Downtime Reason Classification",
            "Defect Type Categorization",
            "Shift Production Logs"
        ], COLOR_AMBER),
        ("04. MEASURE", "Analisis & Optimasi", [
            "Overall Equipment Effectiveness (OEE)",
            "Target vs Actual Gap BI",
            "Downtime Pareto Analytics",
            "Production Quality Trends",
            "Automated Audit-ready Reports"
        ], COLOR_GREEN)
    ]

    card_w = Inches(2.85)
    card_h = Inches(4.95)

    for i, (title, subtitle, items, accent) in enumerate(pillars):
        x = Inches(0.65 + i * 3.06)
        y = Inches(1.70)

        card = create_card(slide, x, y, card_w, card_h, COLOR_BG_WHITE, COLOR_BORDER)
        
        # Header Area
        h_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.75))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = COLOR_BG_CARD
        h_box.line.color.rgb = COLOR_BORDER
        
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.05))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        tb_h = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.10), card_w - Inches(0.3), Inches(0.6))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = 0
        p_th = tf_h.paragraphs[0]
        p_th.text = title
        set_para_font(p_th, FONT_HEADING, 11, COLOR_NAVY, bold=True)
        
        p_subh = tf_h.add_paragraph()
        p_subh.text = subtitle
        set_para_font(p_subh, FONT_BODY, 8.5, COLOR_MUTED_TEXT)

        # List Items
        tb_body = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.90), card_w - Inches(0.3), Inches(3.8))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = 0
        for idx, item in enumerate(items):
            p_item = tf_b.paragraphs[0] if idx == 0 else tf_b.add_paragraph()
            p_item.text = f"•  {item}"
            p_item.space_after = Pt(7)
            set_para_font(p_item, FONT_BODY, 9.5, COLOR_DARK_TEXT)

    add_footer(slide, 4)

# ==========================================
# SLIDE 05: MODULE MAP
# ==========================================
def make_slide_05():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "03. MODULE MAP",
        "Peta Modul Factory-Vision MES",
        "Struktur kapabilitas modular yang dapat diimplementasikan bertahap sesuai kebutuhan operasional pabrik."
    )

    modules = [
        ("PLAN & CONTROL", [
            ("Production Order", "Manajemen pesanan produksi & alokasi kuantitas batch"),
            ("Work Order", "Penerjemahan order ke tugas per stasiun mesin"),
            ("Production Schedule", "Jadwal kerja dan target jam per lini produksi"),
            ("Routing & Process", "Definisi alur tahapan proses manufaktur")
        ], COLOR_BLUE),
        ("SHOP FLOOR EXECUTION", [
            ("Operator Terminal", "UI layar sentuh tablet/PC khusus operator mesin"),
            ("Live Dispatch", "Status berjalan, antrean tugas, dan instruksi lot"),
            ("Machine Monitor", "Status mesin: Running, Idle, Breakdown real-time"),
            ("Shift Control", "Pencatatan shift & serah terima regu kerja")
        ], COLOR_TEAL),
        ("PERFORMANCE & ANALYTICS", [
            ("Executive Dashboard", "KPI makro: OEE, Output, Target vs Actual rate"),
            ("OEE Suite", "Kalkulasi Availability, Performance, Quality"),
            ("Downtime Pareto", "Analisis akar masalah & loss waktu henti mesin"),
            ("Production Reports", "Laporan standar operasional & ekspor format CSV")
        ], COLOR_GREEN),
        ("GOVERNANCE & DATA", [
            ("Quality & Scrap", "Pencatatan defect, scrap rate, dan klasifikasi cacat"),
            ("Batch Traceability", "Silsilah lot dari raw material hingga barang jadi"),
            ("Master Data Engine", "Pusat konfigurasi hierarki pabrik, mesin, dan SKU"),
            ("RBAC & Audit Trail", "Hak akses multi-peran dan log keamanan tak terhapus")
        ], COLOR_NAVY)
    ]

    card_w = Inches(2.85)
    card_h = Inches(4.95)

    for i, (group_name, mod_list, accent) in enumerate(modules):
        x = Inches(0.65 + i * 3.06)
        y = Inches(1.70)

        card = create_card(slide, x, y, card_w, card_h, COLOR_BG_WHITE, COLOR_BORDER)
        
        # Header Pill
        h_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.46))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = accent
        h_box.line.fill.background()
        p_gh = h_box.text_frame.paragraphs[0]
        p_gh.text = group_name
        p_gh.alignment = PP_ALIGN.CENTER
        set_para_font(p_gh, FONT_HEADING, 9.5, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

        for j, (m_title, m_desc) in enumerate(mod_list):
            my = y + Inches(0.58 + j * 1.05)
            m_card = create_card(slide, x + Inches(0.12), my, card_w - Inches(0.24), Inches(0.96), COLOR_BG_CARD, COLOR_BORDER)
            
            tb = slide.shapes.add_textbox(x + Inches(0.18), my + Inches(0.06), card_w - Inches(0.36), Inches(0.82))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = 0
            
            pmt = tf.paragraphs[0]
            pmt.text = m_title
            set_para_font(pmt, FONT_HEADING, 9.5, COLOR_NAVY, bold=True)
            
            pmd = tf.add_paragraph()
            pmd.text = m_desc
            pmd.space_before = Pt(2)
            set_para_font(pmd, FONT_BODY, 8.5, COLOR_MED_TEXT)

    add_footer(slide, 5)

# ==========================================
# MODULE DETAIL TEMPLATE WITH REAL LIGHT SCREENSHOTS
# ==========================================
def make_module_slide_with_image(
    page_num,
    section_label,
    module_title,
    subtitle_text,
    problem_text,
    workflow_steps,
    screenshot_filename,
    screen_title,
    callout_items,
    output_values_list,
    accent_color=COLOR_BLUE
):
    slide = prs.slides.add_slide(blank_layout)

    split_x = Inches(4.5)
    
    # --- RIGHT COLUMN HERO BACKGROUND ---
    hero_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, split_x, Inches(0), Inches(13.333) - split_x, Inches(7.5))
    hero_bg.fill.solid()
    hero_bg.fill.fore_color.rgb = COLOR_NAVY
    hero_bg.line.fill.background()
    
    # --- LEFT COLUMN (White) ---
    left_x = Inches(0.5)
    left_w = Inches(3.7)

    # Custom Header for Left Column
    tx_kicker = slide.shapes.add_textbox(left_x, Inches(0.4), left_w, Inches(0.25))
    tf_k = tx_kicker.text_frame
    tf_k.margin_left = tf_k.margin_top = 0
    pk = tf_k.paragraphs[0]
    pk.text = f"■  {section_label.upper()}"
    set_para_font(pk, FONT_HEADING, 8.5, COLOR_BLUE, bold=True)

    tx_title = slide.shapes.add_textbox(left_x, Inches(0.7), left_w, Inches(1.0))
    tf_t = tx_title.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = 0
    pt = tf_t.paragraphs[0]
    pt.text = module_title
    set_para_font(pt, FONT_HEADING, 18, COLOR_NAVY, bold=True)

    # 1. Problem Card
    p_card = create_card(slide, left_x, Inches(2.0), left_w, Inches(1.36), COLOR_BG_WHITE, COLOR_BORDER)
    strip_p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Inches(2.0), left_w, Inches(0.04))
    strip_p.fill.solid()
    strip_p.fill.fore_color.rgb = COLOR_RED
    strip_p.line.fill.background()

    tb_p = slide.shapes.add_textbox(left_x + Inches(0.15), Inches(2.1), left_w - Inches(0.3), Inches(1.15))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_top = 0
    p1 = tf_p.paragraphs[0]
    p1.text = "MASALAH OPERASIONAL"
    set_para_font(p1, FONT_HEADING, 8.5, COLOR_RED, bold=True)

    p2 = tf_p.add_paragraph()
    p2.text = problem_text
    p2.space_before = Pt(3)
    set_para_font(p2, FONT_BODY, 9.5, COLOR_DARK_TEXT)

    # 2. Workflow Card
    wf_card = create_card(slide, left_x, Inches(3.5), left_w, Inches(1.4), COLOR_BG_WHITE, COLOR_BORDER)
    strip_wf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Inches(3.5), left_w, Inches(0.04))
    strip_wf.fill.solid()
    strip_wf.fill.fore_color.rgb = accent_color
    strip_wf.line.fill.background()

    tb_wf = slide.shapes.add_textbox(left_x + Inches(0.15), Inches(3.6), left_w - Inches(0.3), Inches(1.2))
    tf_wf = tb_wf.text_frame
    tf_wf.word_wrap = True
    tf_wf.margin_left = tf_wf.margin_top = 0
    pw1 = tf_wf.paragraphs[0]
    pw1.text = "MES WORKFLOW"
    set_para_font(pw1, FONT_HEADING, 8.5, accent_color, bold=True)

    pw2 = tf_wf.add_paragraph()
    pw2.text = " → ".join(workflow_steps)
    pw2.space_before = Pt(4)
    set_para_font(pw2, FONT_BODY, 9, COLOR_NAVY, bold=True)

    # 3. Features & Output Card
    val_card = create_card(slide, left_x, Inches(5.05), left_w, Inches(1.8), COLOR_BG_CARD, COLOR_BORDER)
    tb_v = slide.shapes.add_textbox(left_x + Inches(0.15), Inches(5.15), left_w - Inches(0.3), Inches(1.6))
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_top = 0
    pv1 = tf_v.paragraphs[0]
    pv1.text = "OUTPUT & BUSINESS VALUE"
    set_para_font(pv1, FONT_HEADING, 8.5, COLOR_NAVY, bold=True)

    for item in output_values_list[:3]:
        pvi = tf_v.add_paragraph()
        pvi.text = f"✓  {item}"
        pvi.space_before = Pt(3)
        set_para_font(pvi, FONT_BODY, 8.5, COLOR_DARK_TEXT)

    # --- RIGHT COLUMN (Hero Screenshot) ---
    right_x = split_x + Inches(0.4)
    right_w = Inches(13.333) - right_x - Inches(0.4)
    right_h = Inches(5.2)

    # Subtitle acting as Hero Header
    tx_sub = slide.shapes.add_textbox(right_x, Inches(0.5), right_w, Inches(0.6))
    tf_s = tx_sub.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = 0
    ps = tf_s.paragraphs[0]
    ps.text = subtitle_text
    set_para_font(ps, FONT_BODY, 13, COLOR_BG_WHITE, bold=False)

    # Browser Frame Window with Embedded Light Mode Screenshot
    b_y = Inches(1.3)
    b_frame = create_card(slide, right_x, b_y, right_w, right_h, COLOR_BG_WHITE, COLOR_NAVY)
    
    # Browser Bar Header
    b_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, b_y, right_w, Inches(0.36))
    b_top.fill.solid()
    b_top.fill.fore_color.rgb = COLOR_BG_CARD
    b_top.line.fill.background()

    # Window Traffic Light Dots
    for d, c in enumerate([COLOR_RED, COLOR_AMBER, COLOR_GREEN]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Inches(0.12 + d * 0.14), b_y + Inches(0.12), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()

    tb_bt = slide.shapes.add_textbox(right_x + Inches(0.65), b_y + Inches(0.06), right_w - Inches(0.75), Inches(0.25))
    tf_bt = tb_bt.text_frame
    tf_bt.margin_left = tf_bt.margin_top = 0
    pbt = tf_bt.paragraphs[0]
    pbt.text = f"Factory-Vision MES  |  {screen_title}"
    set_para_font(pbt, FONT_BODY, 8.5, COLOR_NAVY)

    # Embed Actual Light Screenshot
    img_path = os.path.join(SCREENSHOTS_DIR, screenshot_filename)
    if os.path.exists(img_path):
        img_top = b_y + Inches(0.36)
        img_height = right_h - Inches(0.36)
        slide.shapes.add_picture(img_path, right_x + Inches(0.06), img_top, width=right_w - Inches(0.12), height=img_height)

    # Callout Strip overlapping bottom of right column
    co_card = create_card(slide, right_x, Inches(6.7), right_w, Inches(0.5), COLOR_BG_CARD, COLOR_BLUE)
    tb_co = slide.shapes.add_textbox(right_x + Inches(0.15), Inches(6.75), right_w - Inches(0.3), Inches(0.4))
    tf_co = tb_co.text_frame
    tf_co.word_wrap = True
    tf_co.margin_left = tf_co.margin_top = 0
    pco1 = tf_co.paragraphs[0]
    pco1.text = "KAPABILITAS SISTEM:"
    set_para_font(pco1, FONT_HEADING, 7.5, COLOR_BLUE, bold=True)

    pco2 = tf_co.add_paragraph()
    pco2.text = "  |  ".join([f"① {f}" if idx==0 else f"② {f}" if idx==1 else f"③ {f}" if idx==2 else f"④ {f}" for idx, f in enumerate(callout_items[:4])])
    set_para_font(pco2, FONT_BODY, 8.5, COLOR_DARK_TEXT)

    # Custom Footer
    tx_page = slide.shapes.add_textbox(Inches(12.5), Inches(7.15), Inches(0.5), Inches(0.25))
    tf_p = tx_page.text_frame
    tf_p.margin_left = tf_p.margin_top = 0
    pp = tf_p.paragraphs[0]
    pp.text = f"{page_num:02d}"
    pp.alignment = PP_ALIGN.RIGHT
    set_para_font(pp, FONT_BODY, 9, COLOR_BG_CARD, bold=True)

    tx_brand = slide.shapes.add_textbox(left_x, Inches(7.15), left_w, Inches(0.25))
    tf_b = tx_brand.text_frame
    tf_b.margin_left = tf_b.margin_top = 0
    pb = tf_b.paragraphs[0]
    pb.text = "FACTORY-VISION  |  Manufacturing Execution System"
    set_para_font(pb, FONT_BODY, 8.5, COLOR_MUTED_TEXT)

# ==========================================
# SLIDE 20: END-TO-END PRODUCTION JOURNEY
# ==========================================
def make_slide_20():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "18. END-TO-END JOURNEY",
        "Satu Pesanan Produksi, Dari Perencanaan Hingga Review Eksekutif",
        "Satu siklus terpadu yang menyatukan seluruh modul Factory-Vision dalam satu alur kerja mulus tanpa friksi."
    )

    journey_steps = [
        ("01", "PO Created", "PPIC membuat order PO-088 (25,000 pcs)."),
        ("02", "WO Split", "Sistem memecah jadi WO Mixing & Curing."),
        ("03", "Machine Ready", "WO dialokasikan ke Line Alpha & TBM-001."),
        ("04", "Operator Start", "Operator login di terminal & menekan Start."),
        ("05", "Data Capture", "Output good, reject, & downtime terekam."),
        ("06", "WO Complete", "Target selesai & lot otomatis terdaftar."),
        ("07", "OEE Calculated", "Sistem menghitung efisiensi OEE otomatis."),
        ("08", "Executive BI", "Direksi melihat capaian di dashboard.")
    ]

    card_w = Inches(1.36)
    card_h = Inches(4.35)

    for i, (num, title, desc) in enumerate(journey_steps):
        x = Inches(0.65 + i * 1.52)
        y = Inches(1.70)

        card = create_card(slide, x, y, card_w, card_h, COLOR_BG_WHITE, COLOR_BORDER)
        
        step_col = COLOR_BLUE if i < 3 else (COLOR_TEAL if i < 6 else COLOR_GREEN)
        nb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.55))
        nb.fill.solid()
        nb.fill.fore_color.rgb = step_col
        nb.line.fill.background()
        pnb = nb.text_frame.paragraphs[0]
        pnb.text = f"STEP {num}"
        pnb.alignment = PP_ALIGN.CENTER
        set_para_font(pnb, FONT_HEADING, 9, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

        tb = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.68), card_w - Inches(0.16), Inches(3.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        
        pt = tf.paragraphs[0]
        pt.text = title
        set_para_font(pt, FONT_HEADING, 10, COLOR_NAVY, bold=True)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.space_before = Pt(4)
        set_para_font(pd, FONT_BODY, 8.5, COLOR_MED_TEXT)

    tb_bottom = create_card(slide, Inches(0.65), Inches(6.20), Inches(12.033), Inches(0.68), COLOR_BG_CARD, COLOR_BLUE)
    tb_b = slide.shapes.add_textbox(Inches(0.85), Inches(6.28), Inches(11.6), Inches(0.5))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = 0
    p_b = tf_b.paragraphs[0]
    p_b.text = "Prinsip Inti: Satu aktivitas di lantai produksi secara otomatis mengalirkan data bernilai tinggi ke seluruh level organisasi tanpa input ganda."
    set_para_font(p_b, FONT_BODY, 10, COLOR_NAVY, bold=True)

    add_footer(slide, 20)

# ==========================================
# SLIDE 21: USER JOURNEY SWIMLANE
# ==========================================
def make_slide_21():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "19. USER JOURNEY",
        "Siapa yang Menggunakan Factory-Vision MES?",
        "Setiap peran organisasi mendapatkan antarmuka dan informasi spesifik sesuai tanggung jawab operasionalnya."
    )

    roles = [
        ("PPIC / Planner", "Perencanaan & Jadwal", "Membuat Production Order, menentukan prioritas rilis, memecah Work Order, dan mengalokasikan target kuantitas.", COLOR_BLUE),
        ("Supervisor Lini", "Pengawasan Operasional", "Memantau Live Production, mengalokasikan operator dan mesin, serta merespons kendala downtime mesin secara instan.", COLOR_TEAL),
        ("Operator Mesin", "Eksekusi Shop Floor", "Menjalankan Work Order via Operator Terminal, mencatat output good/reject, serta mencatat alasan jika mesin berhenti.", COLOR_AMBER),
        ("MES Engine", "Pemrosesan Otomatis", "Mengalkulasi Availability, Performance, Quality, OEE, mendeteksi deviasi bottleneck, dan memicu notifikasi peringatan.", COLOR_GREEN),
        ("Plant Manager", "Analisis & Optimasi", "Menganalisis Downtime Pareto, gap Target vs Actual, mengevaluasi produktivitas shift, dan memimpin continuous improvement.", COLOR_NAVY),
        ("Direksi / Eksekutif", "Strategi Bisnis", "Memantau Executive Cockpit makro, utilisasi kapasitas terpasang, tren efisiensi biaya, dan ROI investasi pabrik.", COLOR_NAVY)
    ]

    card_h = Inches(0.70)
    for i, (role_name, role_focus, role_desc, accent) in enumerate(roles):
        y = Inches(1.70 + i * 0.78)
        
        lane = create_card(slide, Inches(0.65), y, Inches(12.033), card_h, COLOR_BG_WHITE, COLOR_BORDER)
        
        pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), y, Inches(2.6), card_h)
        pill.fill.solid()
        pill.fill.fore_color.rgb = accent
        pill.line.fill.background()
        
        tb_pill = slide.shapes.add_textbox(Inches(0.75), y + Inches(0.10), Inches(2.4), Inches(0.5))
        tf_p = tb_pill.text_frame
        tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_top = 0
        pp1 = tf_p.paragraphs[0]
        pp1.text = role_name
        set_para_font(pp1, FONT_HEADING, 10, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        
        pp2 = tf_p.add_paragraph()
        pp2.text = role_focus
        set_para_font(pp2, FONT_BODY, 8, COLOR_BORDER)

        tb_desc = slide.shapes.add_textbox(Inches(3.4), y + Inches(0.14), Inches(9.1), Inches(0.45))
        tf_d = tb_desc.text_frame
        tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_top = 0
        pd = tf_d.paragraphs[0]
        pd.text = role_desc
        set_para_font(pd, FONT_BODY, 9.5, COLOR_DARK_TEXT)

    add_footer(slide, 21)

# ==========================================
# SLIDE 22: WHAT CHANGES AFTER MES
# ==========================================
def make_slide_22():
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide,
        "20. BUSINESS VALUE",
        "Transformasi Operasional Sebelum & Sesudah Factory-Vision",
        "Dampak nyata implementasi MES terhadap produktivitas, efisiensi waktu, dan keandalan pabrik Anda."
    )

    rows = [
        ("Pencatatan Data", "Manual pada kertas form fisik, rawan hilang & salah rekap", "Digital real-time via Operator Terminal langsung di mesin"),
        ("Visibilitas Produksi", "Terlambat (diketahui akhir shift atau esok hari)", "Live real-time detik per detik pada monitor kendali"),
        ("Analisis Downtime", "Hanya perkiraan kasar tanpa rincian reason code", "Pareto analitik terperinci: durasi, frekuensi, dan akar masalah"),
        ("Kalkulasi OEE", "Dihitung manual mingguan dengan rumus bervariasi", "OEE terstandarisasi global dikalkulasi otomatis oleh sistem"),
        ("Pelacakan Kualitas", "Scrap tercatat agregat tanpa konteks nomor batch", "Kualitas terhubung langsung ke Work Order, operator, & mesin"),
        ("Penelusuran Recall", "Membutuhkan waktu 2-4 hari mencari berkas gudang", "Silsilah batch lengkap dapat ditelusuri dalam 1 klik (< 1 menit)"),
        ("Pelaporan Manajemen", "Staf lembur membuat rekap harian di spreadsheet", "Laporan otomatis siap pakai dan data dapat diekspor instan")
    ]

    th_y = Inches(1.70)
    th_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), th_y, Inches(12.033), Inches(0.40))
    th_box.fill.solid()
    th_box.fill.fore_color.rgb = COLOR_NAVY
    th_box.line.fill.background()

    tb_th = slide.shapes.add_textbox(Inches(0.8), th_y + Inches(0.06), Inches(11.8), Inches(0.3))
    tf_th = tb_th.text_frame
    tf_th.word_wrap = True
    tf_th.margin_left = tf_th.margin_top = 0
    
    pth = tf_th.paragraphs[0]
    pth.text = f"{'DIMENSI OPERASIONAL':<28} | {'SEBELUM MES (TRADISIONAL)':<45} | {'DENGAN FACTORY-VISION MES'}"
    set_para_font(pth, FONT_HEADING, 9.5, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    for i, (dim, before, after) in enumerate(rows):
        ry = Inches(2.14 + i * 0.53)
        bg = COLOR_BG_CARD if i % 2 == 0 else COLOR_BG_WHITE
        rcard = create_card(slide, Inches(0.65), ry, Inches(12.033), Inches(0.48), bg, COLOR_BORDER)

        tb_r = slide.shapes.add_textbox(Inches(0.8), ry + Inches(0.08), Inches(11.8), Inches(0.35))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        tf_r.margin_left = tf_r.margin_top = 0
        
        pr = tf_r.paragraphs[0]
        pr.text = f"•  {dim:<20} :   [Sebelum] {before}   →   [Factory-Vision] {after}"
        set_para_font(pr, FONT_BODY, 8.5, COLOR_DARK_TEXT)

    closing_box = create_card(slide, Inches(0.65), Inches(5.95), Inches(12.033), Inches(0.85), COLOR_NAVY, COLOR_BLUE)
    tb_c = slide.shapes.add_textbox(Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.65))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = 0
    pc1 = tf_c.paragraphs[0]
    pc1.text = "VISIBILITY  →  CONTROL  →  IMPROVEMENT"
    pc1.alignment = PP_ALIGN.CENTER
    set_para_font(pc1, FONT_HEADING, 15, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    pc2 = tf_c.add_paragraph()
    pc2.text = "Factory-Vision MES: Fondasi Utama Transformasi Pabrik Pintar Menuju Industri 4.0."
    pc2.alignment = PP_ALIGN.CENTER
    set_para_font(pc2, FONT_BODY, 9.5, RGBColor(0x90, 0xCD, 0xF4))

    add_footer(slide, 22)

# ==========================================
# BUILD COMPLETE PRESENTATION DECK
# ==========================================
def build_complete_deck(output_file):
    print("Building Slide 01: Cover...")
    make_slide_01()
    print("Building Slide 02: Challenges...")
    make_slide_02()
    print("Building Slide 03: How MES Connects...")
    make_slide_03()
    print("Building Slide 04: 4 Pillars...")
    make_slide_04()
    print("Building Slide 05: Module Map...")
    make_slide_05()

    print("Building Slide 06: Production Management (with light screenshot)...")
    make_module_slide_with_image(
        6,
        "04. PRODUCTION MANAGEMENT",
        "Production & Work Order Management",
        "Menerjemahkan production plan menjadi instruksi kerja terstruktur yang siap dijalankan operator dan mesin.",
        "Rencana produksi sulit dipecah menjadi tugas harian yang teralokasi jelas ke mesin, routing, dan kapasitas operator.",
        ["Production Plan", "Production Order", "Work Order Breakdown", "Machine Allocation", "Release to Shop Floor"],
        "05-work-orders.png",
        "http://localhost:3100/work-orders",
        ["Production Order Tracking", "Work Order State Machine", "Line & Machine Dispatch", "Real-time Quantity Progress"],
        [
            "Pekerjaan produksi terstruktur dan siap dieksekusi tanpa kebingungan prioritas.",
            "Visibilitas status progress order secara real-time dari level agregat hingga detail.",
            "Pengalokasian kapasitas mesin dan lini menjadi lebih optimal dan terencana."
        ],
        COLOR_BLUE
    )

    print("Building Slide 07: Shop Floor Execution (with light screenshot)...")
    make_module_slide_with_image(
        7,
        "05. SHOP FLOOR EXECUTION",
        "Operator Terminal & Shop Floor Execution",
        "Antarmuka sentuh ramah operator untuk eksekusi perintah kerja dan pencatatan output langsung di stasiun mesin.",
        "Operator terbebani menulis logsheet kertas fisik yang rawan kotor, hilang, dan lambat direkap oleh admin.",
        ["My Work Order", "Start Production", "Run Machine", "Record Output", "Log Downtime", "Complete WO"],
        "20-operator-terminal.png",
        "http://localhost:3200 (Operator Terminal)",
        ["Touch-friendly UI", "Good & Reject Quick Counters", "Instant Downtime Logging", "15-min Session Lock"],
        [
            "Data aktual produksi berasal langsung dari sentuhan pertama di shop floor.",
            "Menghilangkan beban rekap manual di akhir shift bagi operator dan supervisor.",
            "Mencegah jeda waktu update progress ke seluruh level manajemen."
        ],
        COLOR_TEAL
    )

    print("Building Slide 08: Live Production (with light screenshot)...")
    make_module_slide_with_image(
        8,
        "06. LIVE PRODUCTION",
        "Live Production & Shop Floor Monitoring",
        "Pemantauan visual seluruh lini produksi secara real-time untuk mendeteksi bottleneck dan deviasi secara dini.",
        "Supervisor harus berjalan keliling pabrik secara fisik untuk mengetahui mesin mana yang sedang macet atau antre material.",
        ["Live Telemetry", "Status Mesin / WO", "Progress Detection", "Bottleneck Isolation", "Rapid Recovery"],
        "06-live-board.png",
        "http://localhost:3100/live-board",
        ["Machine Status Matrix", "Active Work Order Progress", "Real-time Bottleneck Alert", "Estimated Completion"],
        [
            "Supervisor dapat langsung mengidentifikasi stasiun yang membutuhkan asistensi.",
            "Respon terhadap mesin berhenti menjadi jauh lebih cepat dan terukur.",
            "Transparansi total kondisi shop floor tanpa perlu menunggu laporan lisan."
        ],
        COLOR_BLUE
    )

    print("Building Slide 09: Downtime Management (with light screenshot)...")
    make_module_slide_with_image(
        9,
        "07. DOWNTIME MANAGEMENT",
        "Downtime Management & Loss Analysis",
        "Mengubah mesin berhenti dari sekadar insiden menjadi analisis kehilangan waktu terstruktur untuk perbaikan berkelanjutan.",
        "Penyebab mesin berhenti tidak tercatat dengan klasifikasi jelas, sehingga tindakan perbaikan hanya bersifat reaktif.",
        ["Machine Stop", "Downtime Event Created", "Reason Code Tagging", "Duration Tracking", "Pareto Loss Ranking"],
        "04-downtime.png",
        "http://localhost:3100/downtime-analytics",
        ["Reason Code Categorization", "Planned vs Unplanned Split", "Pareto Impact Analysis", "Machine Downtime History"],
        [
            "Tim maintenance dan produksi mengetahui pasti 20% penyebab yang menimbulkan 80% downtime.",
            "Perbaikan mesin menjadi terarah berdasarkan data durasi dan frekuensi aktual.",
            "Meningkatkan Availability mesin secara konsisten dari waktu ke waktu."
        ],
        COLOR_RED
    )

    print("Building Slide 10: OEE & Performance (with light screenshot)...")
    make_module_slide_with_image(
        10,
        "08. OEE & PERFORMANCE",
        "Overall Equipment Effectiveness (OEE)",
        "Mengukur efektivitas pemanfaatan aset pabrik secara objektif melalui standar global OEE.",
        "Output saja belum cukup; manajemen tidak mengetahui apakah mesin berjalan lambat, sering mati, atau banyak menghasilkan reject.",
        ["Production Data", "Availability Calc", "Performance Calc", "Quality Calc", "OEE Score & Loss Analysis"],
        "03-oee.png",
        "http://localhost:3100/oee",
        ["Standard OEE Formula", "Decomposed Loss Factors", "Trend over Shifts", "Line & Machine Benchmarking"],
        [
            "Standardisasi metrik produktivitas yang diakui secara global bagi manajemen.",
            "Mengetahui letak inefisiensi: apakah karena mesin mati, kecepatan lambat, atau scrap.",
            "Dasar pengambilan keputusan investasi penambahan kapasitas mesin baru."
        ],
        COLOR_GREEN
    )

    print("Building Slide 11: Target vs Actual (with light screenshot)...")
    make_module_slide_with_image(
        11,
        "09. PRODUCTION PERFORMANCE",
        "Target vs Actual & Production Gap Analysis",
        "Evaluasi pencapaian target produksi per shift, per lini, dan per order untuk mencegah keterlambatan kirim.",
        "Kesenjangan rencana dan aktual produksi baru disadari saat tenggat waktu pengiriman pesanan klien sudah mendekat.",
        ["Plan Target", "Actual Output", "Achievement % Calc", "Gap Identification", "Corrective Dispatch"],
        "02-production-performance.png",
        "http://localhost:3100/target-vs-actual",
        ["Real-time Target Tracking", "Achievement % Calculation", "At-Risk Order Warning", "Hourly Pace Analysis"],
        [
            "Production Manager dapat melakukan penyesuaian shift sebelum order mengalami keterlambatan.",
            "Komitmen pengiriman ke pelanggan menjadi jauh lebih terjamin dan terprediksi.",
            "Mengevaluasi akurasi perencanaan kapasitas produksi PPIC secara berkelanjutan."
        ],
        COLOR_BLUE
    )

    print("Building Slide 12: Executive Dashboard (with light screenshot)...")
    make_module_slide_with_image(
        12,
        "10. EXECUTIVE DASHBOARD",
        "Executive Dashboard & Central Command",
        "Satu layar kendali eksekutif untuk memantau performa bisnis pabrik dan mendeteksi anomali operasional.",
        "Eksekutif dan direksi membutuhkan rangkuman performa lintas departemen tanpa harus membaca puluhan lembar laporan terpisah.",
        ["Top Macro KPIs", "Identify Anomaly", "Drill-Down to Line", "Root Cause Analysis", "Strategic Intervention"],
        "01-executive-dashboard.png",
        "http://localhost:3100/ (Executive Dashboard)",
        ["8 Executive KPIs", "Single Screen Operations Cockpit", "Multi-Plant Drilldown", "Scrap & Downtime Financial Loss"],
        [
            "Direksi dan GM pabrik memiliki 'cockpit' kendali real-time 24/7.",
            "Rapat koordinasi harian menjadi jauh lebih singkat, fokus, dan berbasis fakta data.",
            "Kecepatan respon organisasi manufaktur meningkat drastis terhadap fluktuasi pasar."
        ],
        COLOR_NAVY
    )

    print("Building Slide 13: Quality Visibility (with light screenshot)...")
    make_module_slide_with_image(
        13,
        "11. QUALITY VISIBILITY",
        "Quality Visibility & Bottleneck Tracking",
        "Mencatat reject langsung pada konteks proses produksi untuk mengidentifikasi kontributor cacat terbesar.",
        "Data reject sering kali hanya dicatat sebagai angka agregat tanpa keterkaitan dengan nomor batch, mesin, atau operator terkait.",
        ["Production Run", "Capture Good/Reject", "Classify Defect Type", "Link to Batch & Machine", "Bottleneck Pareto"],
        "08-bottlenecks.png",
        "http://localhost:3100/bottlenecks",
        ["Instant Reject Recording", "Defect Pareto Classification", "Lot-based Quality Context", "Scrap Rate Trend"],
        [
            "Tim Quality Assurance dapat langsung mengisolasi lot yang berpotensi cacat.",
            "Akar masalah kualitas (apakah mesin, bahan baku, atau operator) cepat teridentifikasi.",
            "Menurunkan biaya kerugian akibat scrap dan komplain retur pelanggan."
        ],
        COLOR_AMBER
    )

    print("Building Slide 14: Traceability (with light screenshot)...")
    make_module_slide_with_image(
        14,
        "12. TRACEABILITY",
        "End-to-End Production Traceability",
        "Silsilah data produksi komprehensif dari bahan baku, stasiun mesin, hingga menjadi finished good siap kirim.",
        "Ketika terjadi keluhan pelanggan, tim pabrik membutuhkan waktu berhari-hari untuk melacak riwayat produksi lot tersebut.",
        ["Finished Product", "Work Order & Batch", "Process Sequence", "Machine & Operator", "Inspection History"],
        "15-master-products.png",
        "http://localhost:3100/settings?tab=products",
        ["Lot & SKU Master History", "Operator & Machine Log", "QC Inspection Linkage", "One-Click Audit History"],
        [
            "Waktu investigasi recall atau issue kualitas terpangkas dari hari menjadi hitungan menit.",
            "Memenuhi persyaratan audit ISO 9001, GMP, BPOM, dan standar industri ketat lainnya.",
            "Meningkatkan kepercayaan dan kredibilitas di mata audit eksternal dan prinsipal."
        ],
        COLOR_TEAL
    )

    print("Building Slide 15: Shift Management (with light screenshot)...")
    make_module_slide_with_image(
        15,
        "13. SHIFT MANAGEMENT",
        "Shift Management & Handover Continuity",
        "Menjaga kontinuitas operasional antar-shift dengan rekap capaian dan pencatatan isu terbuka secara digital.",
        "Informasi kendala mesin atau sisa target sering terputus saat pergantian shift, memicu hilangnya produktivitas awal shift.",
        ["Shift Execution", "Production Summary", "Open Issues & Downtime", "Digital Handover", "Shift Continuation"],
        "07-shift-handover.png",
        "http://localhost:3100/shift-handover",
        ["Automated Shift Output Log", "Open Machine Notes & Issues", "Remaining Target Transfer", "Two-way Digital Sign-off"],
        [
            "Menghilangkan 'blind spot' dan miskomunikasi saat pergantian jam kerja operator.",
            "Shift berikutnya dapat langsung bekerja tanpa membuang waktu 30-45 menit mencari informasi.",
            "Evaluasi produktivitas dan kedisiplinan kerja per regu shift menjadi transparan dan adil."
        ],
        COLOR_BLUE
    )

    print("Building Slide 16: Reports (with light screenshot)...")
    make_module_slide_with_image(
        16,
        "14. PRODUCTION REPORTS",
        "Standardized Production Reports & Export",
        "Standardisasi laporan operasional pabrik siap audit dan ekspor data komprehensif dalam satu klik.",
        "Staf administrasi menghabiskan waktu berjam-jam setiap hari hanya untuk membuat laporan harian dan mingguan secara manual.",
        ["Select Report Type", "Filter Period & Scope", "Dynamic Aggregation", "Analyze Trend", "Export CSV / PDF"],
        "09-reports-production.png",
        "http://localhost:3100/reports?tab=production",
        ["Automated Pre-built Reports", "Date & Line Range Filtering", "1-Click CSV / Excel Export", "Audit-Ready Summaries"],
        [
            "Menghemat puluhan jam kerja staf administrasi pabrik setiap bulannya.",
            "Menghilangkan risiko kesalahan rumus dan manipulasi data pada laporan spreadsheet manual.",
            "Laporan review manajemen bulanan selalu siap tepat waktu tanpa hambatan."
        ],
        COLOR_GREEN
    )

    print("Building Slide 17: Master Data (with light screenshot)...")
    make_module_slide_with_image(
        17,
        "15. MASTER DATA",
        "Master Data Architecture & Plant Hierarchy",
        "Fondasi data tunggal yang terstandarisasi untuk memastikan integritas seluruh workflow manufaktur.",
        "Perbedaan penamaan kode produk, mesin, dan unit di berbagai departemen memicu kekacauan data produksi.",
        ["Plant Definition", "Line & Machine Master", "Product & Process Routing", "Shift & Reason Codes", "System Consistency"],
        "13-master-machines.png",
        "http://localhost:3100/settings?tab=machines",
        ["Plant & Machine Hierarchy", "Standard Ideal Cycle Times", "Reason Code Repository", "Routing Sequence Config"],
        [
            "Satu sumber kebenaran data (Single Source of Truth) untuk seluruh lini operasi.",
            "Kemudahan penambahan mesin atau lini baru (scalable) seiring ekspansi pabrik.",
            "Memastikan kalkulasi OEE dan target produksi selalu akurat dan konsisten."
        ],
        COLOR_NAVY
    )

    print("Building Slide 18: User & Access (with light screenshot)...")
    make_module_slide_with_image(
        18,
        "16. USER & ACCESS",
        "Role-Based Access Control & User Security",
        "Tata kelola hak akses berbasis peran untuk menjamin keamanan data dan integritas operasional pabrik.",
        "Operator atau staf yang tidak berwenang dapat secara tidak sengaja mengubah parameter atau jadwal pesanan penting.",
        ["User Directory", "Role Assignment", "Granular Permissions", "Scope Enforcement", "Secure Actions"],
        "17-master-roles.png",
        "http://localhost:3100/settings?tab=roles",
        ["Role-Based Permissions Matrix", "Dedicated Terminal Scope", "Multi-User Support", "Granular Action Limits"],
        [
            "Setiap staf hanya fokus pada antarmuka dan fungsi yang relevan dengan tugasnya.",
            "Mencegah kesalahan operasional dan perubahan data krusial tanpa otorisasi.",
            "Memenuhi standar keamanan IT enterprise dan tata kelola korporat yang ketat."
        ],
        COLOR_BLUE
    )

    print("Building Slide 19: Audit Trail (with light screenshot)...")
    make_module_slide_with_image(
        19,
        "17. AUDIT TRAIL",
        "Audit Trail & Operational Governance",
        "Rekam jejak digital yang tidak dapat diubah untuk setiap transaksi produksi, intervensi, dan status perubahan.",
        "Ketika terjadi perubahan status order atau koreksi data penting, tidak ada bukti siapa yang melakukannya dan kapan waktu tepatnya.",
        ["User Action", "System Event Capture", "Timestamp & Delta", "Immutable Log", "Audit History Review"],
        "12-audit-logs.png",
        "http://localhost:3100/audit-logs",
        ["Complete Event Logging", "User & IP Traceability", "Before/After Change Delta", "Exportable Audit Records"],
        [
            "Transparansi dan akuntabilitas penuh pada setiap jenjang operasional pabrik.",
            "Mempermudah proses audit internal maupun audit sertifikasi eksternal.",
            "Memberikan ketenangan bagi manajemen puncak atas integritas data sistem."
        ],
        COLOR_NAVY
    )

    print("Building Slide 20: End-to-End Journey...")
    make_slide_20()
    print("Building Slide 21: User Journey Swimlane...")
    make_slide_21()
    print("Building Slide 22: Business Value Closing...")
    make_slide_22()

    prs.save(output_file)
    print(f"Perfected PowerPoint Deck successfully saved to: {output_file}")

if __name__ == "__main__":
    out_file = os.path.join(os.path.dirname(__file__), "..", "Docs", "Factory-Vision-MES-Product-Deck.pptx")
    out_file = os.path.abspath(out_file)
    build_complete_deck(out_file)
