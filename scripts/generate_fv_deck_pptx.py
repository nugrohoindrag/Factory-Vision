"""
Factory Vision MES — Product Deck, PowerPoint build.

Reproduces the published HTML deck as a native 16:9 .pptx, using the product's
own design system rather than a generic consulting template: the navy #0A4174
accent from packages/ui/src/fv/palette.css, the white sheet the sign-in screen
is built on, and the KPI tile with a tone rail on its leading edge.

Geometry is ported 1:1 from the HTML deck, which is authored at 1280x720. At
13.333in x 7.5in that is exactly 96 px per inch, so `px()` below converts
coordinates directly and the two decks stay in step.

    python scripts/generate_fv_deck_pptx.py

Output: Docs/Factory-Vision-MES-Product-Deck.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "Docs",
                                   "Factory-Vision-MES-Product-Deck.pptx"))

# --- Palette, taken verbatim from packages/ui/src/fv/palette.css -------------
NAVY = RGBColor(0x0A, 0x41, 0x74)          # --md-sys-color-primary
NAVY_SOFT = RGBColor(0xBD, 0xD8, 0xE9)     # --primary-container, text on navy
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE_TILE = RGBColor(0xF3, 0xF5, 0xF7)  # --surface-container
TRACK = RGBColor(0xDD, 0xE1, 0xE5)         # --surface-container-highest
OUTLINE = RGBColor(0xDA, 0xDD, 0xE0)       # --outline-variant
ON_SURFACE = RGBColor(0x16, 0x19, 0x1C)
ON_VARIANT = RGBColor(0x56, 0x5C, 0x62)

# Semantic states carry meaning on a shop floor; they are not the accent.
SUCCESS = RGBColor(0x2F, 0xA9, 0x68)
WARNING = RGBColor(0x9A, 0x5B, 0x12)
ERROR = RGBColor(0xB3, 0x26, 0x1E)
INFO = RGBColor(0x2A, 0x6E, 0x96)

# Roboto Flex and Inter are the product's faces but are not installed with
# Office, and PowerPoint substitutes silently rather than falling back the way
# a browser does. Segoe UI is the nearest face every Windows install has.
FONT = "Segoe UI"

PX = 96.0


def px(v):
    """HTML deck pixels -> PowerPoint length."""
    return Inches(v / PX)


prs = Presentation()
prs.slide_width = px(1280)
prs.slide_height = px(720)
BLANK = prs.slide_layouts[6]


# --- primitives --------------------------------------------------------------

def _spacing(run, hundredths):
    """Letter-spacing, which python-pptx does not expose."""
    run.font._rPr.set('spc', str(int(hundredths)))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def write(tf, text, size, color=ON_SURFACE, bold=False, spacing=None,
          align=PP_ALIGN.LEFT, line=None, first=False, caps=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if line:
        p.line_spacing = line
    run = p.add_run()
    run.text = text.upper() if caps else text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if spacing:
        _spacing(run, spacing)
    return p


def rect(slide, x, y, w, h, fill=None, line=None, radius=None, line_w=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    if radius:
        # `adjustment` is a fraction of the shorter side; convert from a pixel
        # radius so a tall card and a short one share the same corner.
        s.adjustments[0] = min(0.5, radius / min(w, h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


# --- composed components -----------------------------------------------------

def eyebrow(slide, text, y=46, color=NAVY, x=56, w=1168):
    """Small uppercase kicker with a hairline running off to the right."""
    tf = textbox(slide, x, y, 420, 16)
    write(tf, text, 8, color, bold=True, spacing=90, first=True, caps=True)
    est = len(text) * 7.3 + 18
    rule = rect(slide, x + est, y + 8, max(60, w - est), 1, fill=OUTLINE)
    return rule


def headline(slide, text, y=74, size=27, w=1168, x=56, color=ON_SURFACE):
    tf = textbox(slide, x, y, w, 90)
    write(tf, text, size, color, bold=True, spacing=-40, line=1.1, first=True)


def lede(slide, text, y, w=800, x=56):
    tf = textbox(slide, x, y, w, 50)
    write(tf, text, 11, ON_VARIANT, line=1.5, first=True)


def tile(slide, x, y, w, h, rail=NAVY, label=None, value=None, caption=None,
         title=None, body=None, body_size=9, value_size=19):
    """The sign-in screen's KPI tile: one flat surface, a hairline, and a rail
    on the leading edge as the only place tone appears."""
    rect(slide, x, y, w, h, fill=SURFACE_TILE, line=OUTLINE, radius=10)
    rect(slide, x, y, 3.5, h, fill=rail, radius=1.5)

    tf = textbox(slide, x + 15, y + 11, w - 27, h - 20)
    first = True
    if label:
        write(tf, label, 7.5, ON_VARIANT, bold=True, spacing=60, first=first, caps=True)
        first = False
    if value:
        p = write(tf, value, value_size, rail, bold=True, spacing=-30, first=first)
        p.space_before = Pt(2)
        first = False
    if title:
        write(tf, title, 10.5, ON_SURFACE, bold=True, spacing=-10, first=first)
        first = False
    if caption:
        p = write(tf, caption, 8.5, ON_VARIANT, line=1.35, first=first)
        p.space_before = Pt(1)
        first = False
    if body:
        p = write(tf, body, body_size, ON_VARIANT, line=1.45, first=first)
        if not first:
            p.space_before = Pt(3)


def panel(slide, x, y, w, h, head=None, sub=None):
    """The navy panel: the one deliberate accent surface on the sheet."""
    rect(slide, x, y, w, h, fill=NAVY, radius=18)
    if head:
        tf = textbox(slide, x + 26, y + 24, w - 52, 90)
        write(tf, head, 14.5, WHITE, bold=True, spacing=-30, line=1.2, first=True)
        if sub:
            p = write(tf, sub, 9.5, NAVY_SOFT, line=1.55)
            p.space_before = Pt(7)


def preview(slide, x, y, w, h, kicker, title):
    """The product preview, anchored to the panel's bottom edge and allowed to
    run past it, the way the sign-in hero does."""
    rect(slide, x, y, w, h, fill=WHITE, radius=18)
    tf = textbox(slide, x + 20, y + 16, w - 40, 40)
    write(tf, kicker, 7, NAVY, bold=True, spacing=90, first=True, caps=True)
    p = write(tf, title, 13, ON_SURFACE, bold=True, spacing=-30)
    p.space_before = Pt(2)


def flow(slide, steps, y=155, x=56, w=1168, seps=None):
    """The inline workflow line under a feature headline. `seps` overrides the
    arrow where the relation between terms is arithmetic rather than sequence."""
    rect(slide, x, y, w, 42, fill=SURFACE_TILE, line=OUTLINE, radius=10)
    tf = textbox(slide, x + 15, y + 13, w - 30, 20)
    p = tf.paragraphs[0]
    for i, step in enumerate(steps):
        if i:
            sep = p.add_run()
            sep.text = "   %s   " % (seps[i - 1] if seps else "→")
            sep.font.name, sep.font.size, sep.font.bold = FONT, Pt(9), True
            sep.font.color.rgb = NAVY
        run = p.add_run()
        run.text = step
        run.font.name, run.font.size, run.font.bold = FONT, Pt(9), True
        run.font.color.rgb = ON_SURFACE


def chain(slide, steps, y, x=56, total_w=1168, cols=4, h=58, gap=30):
    """Process steps on a fixed grid, so rows of three and four stay aligned."""
    cw = (total_w - gap * (cols - 1)) / cols
    for i, step in enumerate(steps):
        col = i % cols
        row = i // cols
        sx = x + col * (cw + gap)
        sy = y + row * (h + 12)
        rect(slide, sx, sy, cw, h, fill=SURFACE_TILE, line=OUTLINE, radius=8)
        rect(slide, sx, sy, 3.5, h, fill=NAVY, radius=1.5)

        tf = textbox(slide, sx + 14, sy + 10, cw - 26, h - 18)
        first = True
        if step.get("num"):
            write(tf, step["num"], 7, NAVY, bold=True, spacing=80, first=True)
            first = False
        write(tf, step["label"], 9.5, ON_SURFACE, bold=True, spacing=-10, first=first)
        if step.get("note"):
            write(tf, step["note"], 7.5, ON_VARIANT, line=1.35)

        # The connector belongs between steps, never after the last of a row.
        last_in_row = (col == cols - 1) or (i == len(steps) - 1)
        if not last_in_row:
            a = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                       px(sx + cw + gap / 2 - 5), px(sy + h / 2 - 5),
                                       px(10), px(10))
            a.rotation = 90
            a.fill.solid()
            a.fill.fore_color.rgb = NAVY
            a.line.fill.background()
            a.shadow.inherit = False


def bar(slide, x, y, w, name, pct, color, name_w=112, val_w=44):
    track_x = x + name_w + 10
    track_w = w - name_w - val_w - 20

    tf = textbox(slide, x, y - 1, name_w, 16)
    write(tf, name, 9, ON_SURFACE, bold=True, first=True)

    rect(slide, track_x, y + 1, track_w, 12, fill=TRACK, radius=6)
    if pct > 0:
        rect(slide, track_x, y + 1, max(14, track_w * pct / 100.0), 12, fill=color, radius=6)

    tf = textbox(slide, x + w - val_w, y - 1, val_w, 16)
    write(tf, f"{pct:.0f}%".replace(".", ","), 9, ON_SURFACE, bold=True,
          align=PP_ALIGN.RIGHT, first=True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def sheet(title, note):
    """A white slide. The sheet is the page; the navy only ever sits on it."""
    slide = prs.slides.add_slide(BLANK)
    bg = rect(slide, 0, 0, 1280, 720, fill=WHITE)
    bg.shadow.inherit = False
    notes(slide, note)
    return slide


# =============================================================================
# Slides
# =============================================================================

def slide_01():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 1280, 720, fill=NAVY)
    notes(s, "Selamat pagi/siang Bapak/Ibu. Hari ini kami memperkenalkan Factory Vision MES, "
             "platform eksekusi manufaktur yang dirancang khusus untuk menjembatani perencanaan di "
             "kantor dengan realitas operasional di lantai pabrik. Dengan tiga pilar utama — Plan, "
             "Execute, dan Measure — pabrik Anda mendapatkan visibilitas dan kendali penuh dari hulu ke hilir.")

    # Brand mark: the FullCircle dot cluster, reduced to its three rings.
    cx, cy = 78, 76
    circle = lambda ox, oy, r, col: rect(s, ox - r, oy - r, r * 2, r * 2, fill=col, radius=r)
    circle(cx, cy, 4.5, WHITE)
    import math
    for k in range(6):
        a = math.radians(k * 60)
        circle(cx + 10 * math.cos(a), cy + 10 * math.sin(a), 3.4, NAVY_SOFT)
    for k in range(12):
        a = math.radians(k * 30)
        circle(cx + 16.5 * math.cos(a), cy + 16.5 * math.sin(a), 2.5, RGBColor(0x7B, 0xBD, 0xE8))

    tf = textbox(s, 104, 62, 600, 40)
    write(tf, "Factory Vision", 15, WHITE, bold=True, spacing=-30, first=True)
    write(tf, "MANUFACTURING EXECUTION SYSTEM  ·  PRODUCT DECK 2026", 7.5, NAVY_SOFT, spacing=60)

    tf = textbox(s, 60, 206, 880, 180)
    write(tf, "Menghubungkan rencana produksi dengan apa yang benar-benar "
              "terjadi di lantai pabrik.", 34, WHITE, bold=True, spacing=-50, line=1.08, first=True)

    tf = textbox(s, 60, 410, 760, 80)
    write(tf, "Satu platform yang menyatukan production planning, shop floor execution, dan "
              "production performance — tanpa rekap kertas, tanpa menunggu akhir shift.",
          12, NAVY_SOFT, line=1.55, first=True)

    pillars = [
        ("PLAN", "Production Order, breakdown Work Order, scheduling, routing, dan alokasi lot."),
        ("EXECUTE", "Terminal sentuh operator, telemetri mesin, kendali WIP dan shift."),
        ("MEASURE", "OEE real-time, Pareto downtime, Target vs Aktual, dan BI eksekutif."),
    ]
    w = (1160 - 24) / 3
    for i, (key, body) in enumerate(pillars):
        x = 60 + i * (w + 12)
        r = rect(s, x, 556, w, 104, fill=RGBColor(0x1B, 0x51, 0x82), radius=10)
        r.line.color.rgb = RGBColor(0x2F, 0x66, 0x96)
        r.line.width = Pt(1)
        tf = textbox(s, x + 18, 574, w - 36, 76)
        write(tf, key, 8.5, NAVY_SOFT, bold=True, spacing=130, first=True)
        p = write(tf, body, 9.5, WHITE, line=1.5)
        p.space_before = Pt(6)


def slide_02():
    s = sheet("Tantangan", "Hampir seluruh pabrik manufaktur menghadapi enam tantangan mendasar ini. "
                           "Masalah utamanya bukan mesin atau orangnya, melainkan ketiadaan sistem digital "
                           "yang menangkap data langsung di detik terjadinya aktivitas produksi.")
    eyebrow(s, "01 · Business Context")
    headline(s, "Apa yang sering menjadi tantangan di operasional produksi?")

    cards = [
        (WARNING, "Data produksi tersebar", "Catatan kertas, spreadsheet, dan percakapan chat menyulitkan rekonsiliasi angka."),
        (WARNING, "Output aktual terlambat", "Kondisi output dan scrap baru diketahui di akhir shift, atau bahkan esok hari."),
        (ERROR, "Pencatatan manual rawan salah", "Operator terbebani mengisi formulir fisik yang mudah keliru dan sulit diverifikasi."),
        (ERROR, "Akar masalah downtime kabur", "Mesin berhenti tanpa alasan terstruktur, sehingga perbaikan sulit diprioritaskan."),
        (WARNING, "Target vs aktual sulit dipantau", "Manajemen tidak punya cara membandingkan gap produksi saat masih bisa ditindak."),
        (ERROR, "Keputusan tertunda", "Respon terhadap hambatan dan bottleneck lambat karena datanya belum tersedia."),
    ]
    w, h, gx, gy = 374, 158, 20, 20
    for i, (rail, title, body) in enumerate(cards):
        x = 56 + (i % 3) * (w + gx)
        y = 210 + (i // 3) * (h + gy)
        tile(s, x, y, w, h, rail=rail, title=title, body=body, body_size=9)


def slide_03():
    s = sheet("Alur MES", "Inilah bagaimana Factory Vision bekerja: menghilangkan jeda waktu informasi antara "
                          "kantor perencanaan PPIC dengan lantai kerja operator secara instan, tanpa rekap fisik berjenjang.")
    eyebrow(s, "02 · MES Overview")
    headline(s, "Factory Vision menghubungkan planning dengan kondisi aktual")
    lede(s, "Tujuh langkah, satu aliran data. Setiap tahap mewariskan konteksnya ke tahap berikutnya, "
            "sehingga tidak ada input yang perlu diketik dua kali.", y=150, w=780)

    chain(s, [
        {"num": "01", "label": "Planning", "note": "PPIC menyusun target dan jadwal."},
        {"num": "02", "label": "Production Order", "note": "Order menjadi batch terstruktur."},
        {"num": "03", "label": "Work Order", "note": "Dipecah per stasiun dan tahapan."},
        {"num": "04", "label": "Shop Floor", "note": "Operator menjalankan instruksi."},
        {"num": "05", "label": "Production Data", "note": "Good, reject, downtime terekam."},
        {"num": "06", "label": "OEE & BI", "note": "Efisiensi dihitung otomatis."},
        {"num": "07", "label": "Decision", "note": "Manajemen bertindak berbasis data."},
    ], y=352, h=64)


def slide_04():
    s = sheet("Empat Pilar", "Factory Vision bukan sekadar dashboard visual, melainkan mesin eksekusi lengkap yang "
                             "memandu operator apa yang harus dikerjakan, mencatat apa yang terjadi, dan mengukur "
                             "hasilnya untuk perbaikan berkelanjutan.")
    eyebrow(s, "02 · MES Overview")
    headline(s, "Satu platform untuk menjalankan, menangkap, dan memahami produksi")

    cols = [
        (NAVY, "Plan", "Production Order\nWork Order decomposition\nScheduling & routing\nLot allocation"),
        (INFO, "Execute", "Terminal sentuh operator\nMachine dispatch\nStart / pause / stop\nShift control & WIP"),
        (SUCCESS, "Capture", "Perhitungan good quantity\nPencatatan reject\nKlasifikasi downtime\nTipe defect"),
        (WARNING, "Measure", "OEE standar global\nGap target vs aktual\nPareto downtime\nLaporan audit"),
    ]
    w = (1168 - 3 * 16) / 4
    for i, (rail, label, body) in enumerate(cols):
        x = 56 + i * (w + 16)
        rect(s, x, 230, w, 300, fill=SURFACE_TILE, line=OUTLINE, radius=10)
        rect(s, x, 230, 3.5, 300, fill=rail, radius=1.5)
        tf = textbox(s, x + 18, 252, w - 34, 260)
        write(tf, label, 9, ON_VARIANT, bold=True, spacing=70, first=True, caps=True)
        for j, ln in enumerate(body.split("\n")):
            p = write(tf, ln, 10, ON_SURFACE, line=1.5)
            p.space_before = Pt(9 if j == 0 else 4)


def slide_05():
    s = sheet("Peta Modul", "Seluruh kapabilitas ini dirancang modular. Klien dapat memulai dari implementasi modul "
                            "inti shop floor terlebih dahulu, lalu bertahap mengaktifkan modul analitik canggih "
                            "sesuai kesiapan tim.")
    eyebrow(s, "03 · Module Map")
    headline(s, "Peta modul Factory Vision MES")
    lede(s, "Modular secara sengaja: mulai dari eksekusi shop floor, aktifkan analitik ketika tim sudah siap.", y=150)

    groups = [
        (NAVY, "Plan & Control", "Production Order  ·  Work Order  ·  Production Schedule  ·  Routing & Process"),
        (INFO, "Shop Floor Execution", "Operator Terminal  ·  Live Dispatch  ·  Machine Monitoring  ·  Shift Control"),
        (WARNING, "Performance & OEE", "Executive Dashboard  ·  OEE Suite  ·  Downtime Pareto  ·  Production Reports"),
        (SUCCESS, "Quality & Governance", "Quality & Scrap  ·  Batch Traceability  ·  Master Data Engine  ·  RBAC & Audit Trail"),
    ]
    w, h = 574, 132
    for i, (rail, title, body) in enumerate(groups):
        x = 56 + (i % 2) * (w + 20)
        y = 240 + (i // 2) * (h + 20)
        tile(s, x, y, w, h, rail=rail, title=title, body=body, body_size=9.5)


def feature(eyebrow_text, title, steps, value, extra, panel_head, panel_sub,
            preview_kicker, preview_title, draw_preview, note):
    """The feature-slide template: narrative on the left, the navy panel with an
    anchored product preview on the right."""
    s = sheet(title, note)
    eyebrow(s, eyebrow_text)
    headline(s, title, size=23)
    flow(s, steps)

    # Left column: the two tiles take their alignment from the panel's own
    # anchors — its heading at the top, its preview at the bottom.
    tile(s, 56, 220, 570, 118, rail=NAVY, label="Nilai bisnis", body=value, body_size=10.5)
    tile(s, 56, 520, 570, 130, rail=extra[0], title=extra[1], body=extra[2], body_size=9)

    panel(s, 660, 220, 564, 430, head=panel_head, sub=panel_sub)
    draw_preview(s)
    return s


def slide_06():
    def prev(s):
        preview(s, 686, 452, 512, 220, "Contoh tampilan", "Production Order · PO-2026-088")
        w = 152
        for i, (rail, label, value, cap) in enumerate([
            (NAVY, "Target", "25.000", "pcs"),
            (SUCCESS, "Released", "3", "work order"),
            (INFO, "Due", "12 Sep", "2026"),
        ]):
            tile(s, 706 + i * (w + 12), 512, w, 74, rail=rail, label=label,
                 value=value, caption=cap, value_size=15)
        for i, (name, pct, col) in enumerate([
            ("Mixing", 100, SUCCESS), ("Filling", 64, NAVY), ("Packaging", 12, WARNING)
        ]):
            bar(s, 706, 604 + i * 22, 472, name, pct, col, name_w=90, val_w=40)

    feature("04 · Production Management",
            "Production & Work Order Management",
            ["Production Plan", "Production Order", "Work Order Breakdown", "Machine Allocation", "Release"],
            "Menjamin setiap tugas terdistribusi ke lini yang tepat, dengan target yang jelas dan dapat "
            "ditagih pertanggungjawabannya.",
            (INFO, "Konteks diwariskan otomatis",
             "Proses, batch/lot, dan mesin melekat pada work order — operator tidak perlu memasukkan ulang."),
            "Rencana menjadi tugas",
            "Satu pesanan pelanggan diturunkan menjadi work order per tahapan, lengkap dengan target dan "
            "lini pelaksananya.",
            "Contoh tampilan", "Production Order · PO-2026-088", prev,
            "Setiap pesanan pelanggan diterjemahkan menjadi Work Order yang jelas pemiliknya. Menjamin setiap "
            "tugas terdistribusi ke lini yang tepat dengan target yang terukur.")


def slide_07():
    def prev(s):
        preview(s, 686, 452, 512, 220, "Contoh tampilan", "Terminal Operator · WO-260829-01")
        w = 152
        for i, (rail, label, value, cap) in enumerate([
            (SUCCESS, "Good", "3.420", "pcs shift ini"),
            (ERROR, "Reject", "28", "0,81% dari total"),
            (INFO, "Cycle", "1,8s", "per unit"),
        ]):
            tile(s, 706 + i * (w + 12), 512, w, 82, rail=rail, label=label,
                 value=value, caption=cap, value_size=17)
        chips = [("Catat Output", NAVY, WHITE), ("Log Downtime", TRACK, ON_VARIANT),
                 ("Selesaikan WO", TRACK, ON_VARIANT)]
        cx = 706
        for text, bg, fg in chips:
            cw = len(text) * 6.2 + 22
            rect(s, cx, 608, cw, 26, fill=bg, radius=13)
            tf = textbox(s, cx, 614, cw, 16)
            write(tf, text, 8, fg, bold=True, align=PP_ALIGN.CENTER, first=True)
            cx += cw + 8

    feature("05 · Shop Floor Execution",
            "Operator Terminal & Shop Floor Execution",
            ["My Work Order", "Start Production", "Run Machine", "Catat Good / Reject", "Log Downtime", "Complete"],
            "Menghilangkan beban rekap kertas sepenuhnya, dan memastikan data valid langsung dari sumber pertama.",
            (SUCCESS, "Tetap bekerja saat Wi-Fi putus",
             "Terminal menyimpan antrean secara lokal dan menyinkronkannya saat koneksi kembali — tanpa duplikasi catatan."),
            "Dirancang untuk tangan bersarung",
            "Tombol besar, angka besar, tanpa formulir panjang. Operator mencatat sambil bekerja, bukan setelahnya.",
            "Contoh tampilan", "Terminal Operator", prev,
            "Terminal operator dirancang untuk tangan yang memakai sarung tangan: tombol besar, angka besar, "
            "dan tidak ada formulir panjang. Menghilangkan seluruh beban rekap kertas dan memastikan data "
            "valid langsung dari sumber pertama.")


def slide_08():
    def prev(s):
        preview(s, 686, 430, 512, 242, "Contoh tampilan", "Live Production Board")
        rows = [(SUCCESS, "Line 1 · Mixing", "Running"),
                (WARNING, "Line 2 · Filling", "Idle · menunggu material"),
                (ERROR, "Line 3 · Packaging", "Breakdown · heater"),
                (SUCCESS, "Line 4 · Labeling", "Running")]
        for i, (rail, name, state) in enumerate(rows):
            y = 492 + i * 44
            rect(s, 706, y, 472, 36, fill=SURFACE_TILE, line=OUTLINE, radius=8)
            rect(s, 706, y, 3.5, 36, fill=rail, radius=1.5)
            rect(s, 722, y + 15, 7, 7, fill=rail, radius=3.5)
            tf = textbox(s, 738, y + 11, 240, 16)
            write(tf, name, 9.5, ON_SURFACE, bold=True, first=True)
            tf = textbox(s, 950, y + 11, 218, 16)
            write(tf, state, 8.5, ON_VARIANT, align=PP_ALIGN.RIGHT, first=True)

    feature("06 · Live Production",
            "Live Production & Shop Floor Monitoring",
            ["Live Grid Telemetry", "Status Mesin", "Progress Detection", "Bottleneck Isolation", "Rapid Recovery"],
            "Memangkas waktu respon supervisor dan teknisi terhadap mesin yang berhenti, dari puluhan menit "
            "menjadi seketika.",
            (WARNING, "Status yang terbaca dari kejauhan",
             "Warna membawa arti operasional, bukan dekorasi — running, idle, dan breakdown terbedakan dalam sekali lihat."),
            "Terbaca dari seberang ruangan",
            "Lini yang berhenti terlihat sebelum ada yang melapor.",
            "Contoh tampilan", "Live Production Board", prev,
            "Supervisor tidak lagi berjalan keliling untuk tahu mesin mana yang berhenti. Memangkas waktu "
            "respon supervisor dan teknisi terhadap mesin berhenti.")


def slide_09():
    def prev(s):
        preview(s, 686, 452, 512, 220, "Contoh tampilan", "Pareto Alasan Downtime · 30 hari")
        for i, (name, pct, col) in enumerate([
            ("Nozzle jammed", 42, ERROR), ("Waiting material", 25, WARNING),
            ("Format changeover", 18, WARNING), ("Sensor trip", 15, INFO),
        ]):
            bar(s, 706, 524 + i * 30, 472, name, pct, col, name_w=118, val_w=40)

    feature("07 · Downtime Management",
            "Downtime Management & Loss Analysis",
            ["Machine Stop", "Downtime Event", "Reason Code Tagging", "Duration Tracking", "Pareto Loss Ranking"],
            "Mengarahkan fokus preventive maintenance pada 20% penyebab yang menimbulkan 80% downtime.",
            (ERROR, "Alasan wajib, bukan opsional",
             "Downtime tanpa reason code adalah angka yang tidak bisa ditindaklanjuti siapa pun."),
            "Dari keluhan menjadi prioritas",
            "Setiap berhentinya mesin wajib punya reason code. Itulah yang mengubah cerita lisan menjadi "
            "peringkat kerugian.",
            "Contoh tampilan", "Pareto Alasan Downtime", prev,
            "Setiap berhentinya mesin wajib punya alasan terstruktur. Itulah yang mengubah keluhan menjadi "
            "Pareto, dan Pareto menjadi prioritas maintenance.")


def slide_10():
    s = sheet("OEE", "OEE dihitung dengan rumus standar global, dari transaksi shop floor yang sama — bukan dari "
                     "rekap terpisah. Standardisasi metrik produktivitas kelas dunia untuk justifikasi investasi "
                     "mesin dan efisiensi aset.")
    eyebrow(s, "08 · OEE & Performance")
    headline(s, "Overall Equipment Effectiveness", size=23)
    flow(s, ["Availability 89,6%", "Performance 94,2%", "Quality 97,6%", "OEE 82,4%"],
         seps=["×", "×", "="])

    w = 182
    for i, (rail, label, value, cap) in enumerate([
        (INFO, "Availability", "89,6%", "run time ÷ planned"),
        (WARNING, "Performance", "94,2%", "vs ideal cycle"),
        (SUCCESS, "Quality", "97,6%", "good ÷ total"),
    ]):
        tile(s, 56 + i * (w + 12), 220, w, 100, rail=rail, label=label, value=value,
             caption=cap, value_size=18)

    tile(s, 56, 340, 570, 120, rail=NAVY, label="Nilai bisnis",
         body="Standardisasi metrik produktivitas kelas dunia, sebagai dasar justifikasi investasi mesin "
              "dan efisiensi aset.", body_size=10.5)
    tile(s, 56, 480, 570, 170, rail=NAVY, title="Satu rumus untuk semua lini",
         body="Availability, Performance, dan Quality dihitung dari transaksi shop floor yang sama, dengan "
              "definisi yang tidak berubah antar orang atau antar minggu.", body_size=9.5)

    panel(s, 660, 220, 564, 430, head="Satu rumus untuk semua lini",
          sub="Perbandingan antar lini menjadi bermakna karena semuanya diukur dengan cara yang sama.")
    preview(s, 686, 452, 512, 220, "Contoh tampilan", "Peringkat OEE per Lini")
    for i, (name, pct, col) in enumerate([
        ("Line Alpha", 88, SUCCESS), ("Line Beta", 75, WARNING),
        ("Line Gamma", 59, ERROR), ("Line Delta", 82, SUCCESS),
    ]):
        bar(s, 706, 524 + i * 30, 472, name, pct, col, name_w=96, val_w=44)


def slide_11():
    def prev(s):
        preview(s, 686, 430, 512, 242, "Contoh tampilan", "Target vs Produksi Aktual")
        tile(s, 706, 492, 230, 76, rail=SUCCESS, label="Pencapaian", value="95,0%",
             caption="142.500 / 150.000 pcs", value_size=17)
        tile(s, 948, 492, 230, 76, rail=ERROR, label="Terancam telat", value="PO-090",
             caption="proyeksi −4.200 unit", value_size=17)
        for i, (name, pct, col) in enumerate([
            ("PO-088", 100, SUCCESS), ("PO-089", 93, WARNING), ("PO-090", 72, ERROR)
        ]):
            bar(s, 706, 588 + i * 26, 472, name, pct, col, name_w=80, val_w=44)

    feature("09 · Production Performance",
            "Target vs Aktual & Production Gap Analysis",
            ["Target", "Actual Counter", "Achievement %", "Gap Identification", "At-Risk Alert"],
            "Menjamin ketepatan waktu pengiriman (On-Time In-Full) ke pelanggan, dengan peringatan sebelum terlambat.",
            (ERROR, "Peringatan dini, bukan laporan kematian",
             "Order yang terancam telat ditandai saat proyeksi menunjukkan gap, bukan saat due date terlewat."),
            "Terlihat saat masih bisa dikejar",
            "Gap ditandai berdasarkan proyeksi akhir periode, bukan setelah tanggal kirim terlewat.",
            "Contoh tampilan", "Target vs Produksi Aktual", prev,
            "Gap produksi terlihat saat masih bisa dikejar, bukan setelah pengiriman terlambat. Menjamin "
            "ketepatan waktu pengiriman On-Time In-Full ke pelanggan.")


def slide_12():
    s = sheet("Executive Dashboard", "Menjadikan rapat direksi fokus pada tindakan berbasis data, bukan perdebatan "
                                     "dokumen. Lima pertanyaan bisnis, dijawab dari transaksi shop floor yang sama.")
    eyebrow(s, "10 · Executive Dashboard")
    headline(s, "Executive Dashboard & Central Command", size=23)
    flow(s, ["5 Pertanyaan Direksi", "Identifikasi Anomali", "Drill-down Lini", "Keputusan Strategis"])

    panel(s, 56, 220, 1168, 430, head="Lima pertanyaan, satu sumber angka",
          sub="Setiap kartu dibaca dari transaksi shop floor. Tidak ada yang dihitung ulang di browser, dan "
              "tidak ada yang perlu direkap semalaman sebelum rapat.")
    preview(s, 92, 430, 1096, 242, "Contoh tampilan", "Executive Dashboard · 7 hari terakhir")
    w = 256
    for i, (rail, label, value, cap) in enumerate([
        (WARNING, "OEE", "82,4%", "Target 85%  ·  −2,6"),
        (SUCCESS, "Pencapaian", "95,0%", "Target 100%  ·  −5,0"),
        (ERROR, "Nilai Scrap", "Rp 14,2jt", "periode berjalan"),
        (NAVY, "Active Orders", "18", "sedang berjalan"),
    ]):
        tile(s, 112 + i * (w + 14), 500, w, 96, rail=rail, label=label, value=value,
             caption=cap, value_size=19)


def slide_13():
    def prev(s):
        preview(s, 686, 452, 512, 220, "Contoh tampilan", "Klasifikasi Cacat")
        for i, (name, pct, col) in enumerate([
            ("Capping seal loose", 46.9, ERROR), ("Underfill", 30.8, WARNING), ("Label scratch", 22.3, INFO)
        ]):
            bar(s, 706, 530 + i * 32, 472, name, pct, col, name_w=126, val_w=46)

    feature("11 · Quality Visibility",
            "Quality Visibility & Defect Tracking",
            ["Output", "Good / Reject Split", "Kategorisasi Defect", "Lot Linkage", "Quality Pareto"],
            "Menurunkan biaya kerugian scrap, dan mempercepat isolasi lot yang berpotensi cacat.",
            (SUCCESS, "Reject rate 2,4%",
             "Di bawah ambang 3,0%. Angkanya dibaca dari catatan produksi, bukan dari rekap mingguan."),
            "Reject yang punya alamat",
            "Setiap reject membawa kategori, lot, dan mesin asalnya — sehingga perbaikannya bisa diarahkan.",
            "Contoh tampilan", "Klasifikasi Cacat", prev,
            "Reject bukan sekadar angka; ia punya kategori, lot, dan mesin asal. Menurunkan biaya kerugian "
            "scrap dan mempercepat isolasi lot berpotensi cacat.")


def slide_14():
    def prev(s):
        preview(s, 686, 424, 512, 248, "Contoh tampilan", "Lot Explorer · LOT-SKN-202608-A1")
        rows = [(NAVY, "Mixing · MIX-001", "28 Agu 07:12  ·  Operator Budi S."),
                (NAVY, "Filling · FIL-002", "28 Agu 09:40  ·  Operator Rina W."),
                (NAVY, "Packaging · PCK-001", "28 Agu 13:05  ·  Operator Dedi P."),
                (SUCCESS, "QC Release", "28 Agu 15:22  ·  Lulus uji kebocoran")]
        for i, (rail, title, cap) in enumerate(rows):
            y = 486 + i * 46
            rect(s, 706, y, 472, 38, fill=SURFACE_TILE, line=OUTLINE, radius=8)
            rect(s, 706, y, 3.5, 38, fill=rail, radius=1.5)
            tf = textbox(s, 722, y + 7, 440, 26)
            write(tf, title, 9, ON_SURFACE, bold=True, first=True)
            write(tf, cap, 7.5, ON_VARIANT)

    feature("12 · Traceability",
            "End-to-End Production Traceability",
            ["Finished Good Lot", "PO & WO Hierarchy", "Station Logs", "Operator Signature", "QC Lab History"],
            "Waktu penelusuran investigasi recall terpangkas dari tiga hari menjadi kurang dari satu menit.",
            (INFO, "Silsilah, bukan arsip",
             "Setiap tahapan menyimpan tautan ke tahap sebelumnya, sehingga penelusuran berjalan otomatis ke hulu."),
            "Silsilah, bukan tumpukan arsip",
            "Ketika ada keluhan pelanggan, pertanyaannya selalu sama: lot ini dibuat kapan, di mesin mana, oleh siapa.",
            "Contoh tampilan", "Lot Explorer", prev,
            "Ketika ada keluhan pelanggan, pertanyaannya selalu sama: lot ini dibuat kapan, di mesin mana, oleh "
            "siapa. Waktu penelusuran investigasi recall terpangkas dari tiga hari menjadi kurang dari satu menit.")


def slide_15():
    def prev(s):
        preview(s, 686, 442, 512, 230, "Contoh tampilan", "Serah Terima Shift")
        tile(s, 706, 504, 230, 78, rail=INFO, label="Shift 1 · keluar", value="4.180",
             caption="good  ·  96 menit henti", value_size=16)
        tile(s, 948, 504, 230, 78, rail=NAVY, label="Shift 2 · masuk", value="2 isu",
             caption="terbuka, perlu tindak lanjut", value_size=16)
        rect(s, 706, 598, 472, 52, fill=SURFACE_TILE, line=OUTLINE, radius=8)
        rect(s, 706, 598, 3.5, 52, fill=WARNING, radius=1.5)
        tf = textbox(s, 722, 610, 442, 36)
        write(tf, "“Heater packaging masih intermiten sejak 13:40. Teknisi sudah dihubungi, "
                  "suku cadang tiba besok pagi.”", 8.5, ON_SURFACE, line=1.4, first=True)

    feature("13 · Shift Management",
            "Shift Management & Handover Continuity",
            ["Shift Summary", "Open Issues & Downtime", "Digital Sign-off", "Kontinuitas Shift Berikutnya"],
            "Meniadakan kehilangan produktivitas di 30 menit awal pergantian regu kerja.",
            (WARNING, "Catatan yang terbawa",
             "Isu terbuka dari shift sebelumnya muncul otomatis di layar supervisor yang masuk."),
            "Regu berikutnya mulai dengan konteks",
            "Hasil shift dan isu yang masih terbuka berpindah bersama tanggung jawabnya.",
            "Contoh tampilan", "Serah Terima Shift", prev,
            "Tiga puluh menit pertama setiap pergantian regu biasanya hilang untuk saling bertanya. "
            "Serah terima digital meniadakan kehilangan produktivitas itu.")


def slide_16():
    s = sheet("Laporan", "Laporan bukan lagi pekerjaan lembur admin. Menghemat puluhan jam kerja administrasi dan "
                         "menjamin laporan bebas manipulasi karena dibaca langsung dari transaksi.")
    eyebrow(s, "14 · Production Reports")
    headline(s, "Standardized Production Reports & Export", size=23)
    flow(s, ["Pilih Laporan", "Filter Periode", "Agregasi Instan", "Ekspor CSV / PDF"])

    reports = [
        (NAVY, "Daily Summary", "Ringkasan produksi harian per lini dan shift."),
        (WARNING, "OEE Trend", "Tren Availability, Performance, dan Quality."),
        (ERROR, "Downtime Pareto", "Peringkat kerugian menurut alasan berhenti."),
        (INFO, "WO Audit Trail", "Jejak lengkap perubahan setiap work order."),
    ]
    w = (1168 - 3 * 16) / 4
    for i, (rail, title, body) in enumerate(reports):
        tile(s, 56 + i * (w + 16), 240, w, 200, rail=rail, title=title, body=body, body_size=9.5)

    tile(s, 56, 480, 1168, 100, rail=NAVY, label="Nilai bisnis",
         body="Menghemat puluhan jam kerja administrasi, dan menjamin laporan bebas manipulasi karena dibaca "
              "langsung dari transaksi.", body_size=11)


def slide_17():
    s = sheet("Master Data", "Master data adalah fondasinya. Menjadi Single Source of Truth yang mudah dikembangkan "
                             "saat pabrik berekspansi ke lini atau lokasi baru.")
    eyebrow(s, "15 · Master Data")
    headline(s, "Master Data Architecture & Plant Hierarchy", size=23)

    chain(s, [
        {"label": "Plant", "note": "Lokasi pabrik"},
        {"label": "Production Line", "note": "Lini produksi"},
        {"label": "Machine Center", "note": "Work center & mesin"},
        {"label": "SKU & BOM", "note": "Produk dan komposisi"},
        {"label": "Routing Sequence", "note": "Urutan proses per produk"},
        {"label": "Reason Codes", "note": "Alasan downtime dan reject"},
        {"label": "Cycle Rates", "note": "Ideal cycle time per produk × mesin"},
    ], y=268, h=68)

    tile(s, 56, 500, 1168, 110, rail=NAVY, label="Nilai bisnis",
         body="Menjadi Single Source of Truth yang mudah dikembangkan saat pabrik berekspansi — tanpa "
              "membangun ulang konfigurasi dari nol.", body_size=11)


def slide_18():
    s = sheet("RBAC", "Setiap peran melihat persis apa yang relevan baginya. Menjaga keamanan sistem dan "
                      "memastikan operator tidak dibebani fungsi yang bukan urusannya.")
    eyebrow(s, "16 · User & Access")
    headline(s, "Role-Based Access Control & User Security", size=23)
    flow(s, ["User Directory", "Assigned Role", "Granular Permission Scope", "Action Control"])

    roles = [
        (NAVY, "Executive", "KPI makro dan tren lintas pabrik. Hanya baca."),
        (INFO, "Production Manager", "Perencanaan, rilis order, dan analisis performa."),
        (SUCCESS, "Supervisor", "Monitoring lini, koreksi data, dan serah terima shift."),
        (WARNING, "Operator Terminal", "Eksekusi work order pada lini yang ditugaskan saja."),
        (INFO, "PPIC & Quality", "Penjadwalan, batch, dan penelusuran mutu."),
        (ERROR, "Administrator", "Master data, pengguna, dan konfigurasi sistem."),
    ]
    w, h = 374, 118
    for i, (rail, title, body) in enumerate(roles):
        x = 56 + (i % 3) * (w + 20)
        y = 226 + (i // 3) * (h + 18)
        tile(s, x, y, w, h, rail=rail, title=title, body=body, body_size=9)

    tile(s, 56, 500, 1168, 76, rail=NAVY,
         body="Izin ditegakkan di server, bukan disembunyikan di antarmuka — sehingga batasannya nyata, "
              "bukan kosmetik.", body_size=10.5)


def slide_19():
    def prev(s):
        preview(s, 686, 442, 512, 230, "Contoh tampilan", "Audit Trail")
        rows = [(WARNING, "CORRECTION_APPROVED · production_record", "30 Agu 14:02  ·  Agung W.  ·  10.4.2.18"),
                (INFO, "WO_STATUS_CHANGED · wo-102", "30 Agu 13:47  ·  Rina W.  ·  terminal TAB-03"),
                (NAVY, "PERMISSION_CHANGED · role Supervisor", "30 Agu 09:15  ·  Administrator  ·  10.4.1.2")]
        for i, (rail, title, cap) in enumerate(rows):
            y = 506 + i * 50
            rect(s, 706, y, 472, 42, fill=SURFACE_TILE, line=OUTLINE, radius=8)
            rect(s, 706, y, 3.5, 42, fill=rail, radius=1.5)
            tf = textbox(s, 722, y + 9, 440, 30)
            write(tf, title, 8.5, ON_SURFACE, bold=True, first=True)
            write(tf, cap, 7.5, ON_VARIANT)

    feature("17 · Audit Trail",
            "Audit Trail & Operational Governance",
            ["User Action", "System Capture", "Timestamp + IP", "Immutable Log"],
            "Akuntabilitas penuh, dan kesiapan mutlak saat menghadapi audit sertifikasi ISO atau BPOM.",
            (ERROR, "Append-only",
             "Catatan tidak dapat diubah maupun dihapus. Koreksi data dicatat sebagai peristiwa baru, "
             "lengkap dengan penyetujunya."),
            "Tercatat, tidak bisa dihapus",
            "Setiap tindakan membawa pelaku, waktu, dan asal terminalnya.",
            "Contoh tampilan", "Audit Trail", prev,
            "Audit trail bersifat append-only: tidak bisa diubah maupun dihapus. Akuntabilitas penuh dan "
            "kemudahan mutlak saat menghadapi audit sertifikasi ISO atau BPOM.")


def slide_20():
    s = sheet("End-to-End Journey", "Satu kali operator menekan tombol di shop floor, data mengalir otomatis "
                                    "melayani kebutuhan operator, supervisor, manajer pabrik, QA, hingga direksi "
                                    "— tanpa input ganda.")
    eyebrow(s, "18 · End-to-End Journey")
    headline(s, "Satu pesanan produksi, dari perencanaan hingga review eksekutif")
    lede(s, "Satu penekanan tombol di lantai pabrik, dan angkanya sudah sampai di rapat direksi. "
            "Tidak ada rekap ulang di antaranya.", y=150, w=820)

    chain(s, [
        {"num": "01", "label": "PO Created"},
        {"num": "02", "label": "WO Split"},
        {"num": "03", "label": "Machine Assigned"},
        {"num": "04", "label": "Operator Start"},
        {"num": "05", "label": "Capture Telemetry"},
        {"num": "06", "label": "WO Completed"},
        {"num": "07", "label": "OEE Calculated"},
        {"num": "08", "label": "Executive Review"},
    ], y=364, h=56)


def slide_21():
    s = sheet("User Journey", "Enam peran, satu sistem. Masing-masing melihat lapisan yang relevan bagi "
                              "pekerjaannya, dari operator di lantai hingga direksi di ruang rapat.")
    eyebrow(s, "19 · User Journey")
    headline(s, "Siapa yang menggunakan Factory Vision MES?", size=23)

    lanes = [
        ("PPIC", "Menyusun perencanaan pesanan, jadwal produksi, dan target batch."),
        ("Supervisor", "Memantau produksi live dan merespon cepat kendala mesin di lini."),
        ("Operator", "Menjalankan work order, mencatat good/reject dan downtime lewat tablet sentuh."),
        ("MES Engine", "Mengotomasi kalkulasi OEE, agregasi data, dan deteksi deviasi."),
        ("Plant Manager", "Menganalisis tren kerugian, mengevaluasi shift, dan continuous improvement."),
        ("Direksi", "Memantau KPI makro bisnis dan efisiensi utilisasi aset pabrik."),
    ]
    for i, (role, body) in enumerate(lanes):
        y = 210 + i * 74
        rect(s, 56, y, 168, 60, fill=NAVY, radius=8)
        tf = textbox(s, 70, y + 21, 148, 20)
        write(tf, role, 10, WHITE, bold=True, spacing=-10, first=True)

        rect(s, 236, y, 988, 60, fill=SURFACE_TILE, line=OUTLINE, radius=8)
        tf = textbox(s, 254, y + 21, 952, 20)
        write(tf, body, 10, ON_VARIANT, first=True)


def slide_22():
    s = sheet("Business Value", "Inilah perubahan konkretnya. Factory Vision MES adalah fondasi utama transformasi "
                                "pabrik pintar menuju Industri 4.0: visibility lebih dulu, lalu control, lalu "
                                "improvement yang berkelanjutan.")
    eyebrow(s, "20 · Business Value")
    headline(s, "Transformasi operasional sebelum & sesudah Factory Vision", size=23)

    rows = [
        ("Pencatatan data", "Manual di kertas, rawan hilang dan salah tulis", "Digital real-time lewat terminal operator"),
        ("Visibilitas lini", "Terlambat, baru diketahui akhir shift atau esok hari", "Live, detik per detik di layar"),
        ("Analisis downtime", "Estimasi kasar tanpa penyebab pasti", "Pareto terperinci dengan durasi dan reason code"),
        ("Kalkulasi OEE", "Manual mingguan, rumusnya berbeda antar orang", "Standar global, dihitung otomatis oleh sistem"),
        ("Penelusuran lot", "Butuh 2–4 hari mencari berkas fisik", "Silsilah lengkap dalam satu klik, di bawah satu menit"),
        ("Pelaporan manajemen", "Admin lembur menyusun rekap spreadsheet", "Laporan otomatis siap pakai dan ekspor instan"),
    ]
    x, y, w = 56, 208, 1168
    col = [250, 440, 478]
    head_h, row_h = 34, 46

    rect(s, x, y, w, head_h, fill=NAVY, radius=6)
    for i, label in enumerate(["Dimensi Operasional", "Sebelum MES", "Dengan Factory Vision"]):
        tf = textbox(s, x + 16 + sum(col[:i]), y + 11, col[i] - 20, 18)
        write(tf, label, 7.5, WHITE, bold=True, spacing=70, first=True, caps=True)

    for r, (dim, before, after) in enumerate(rows):
        ry = y + head_h + r * row_h
        rect(s, x, ry + row_h - 1, w, 1, fill=OUTLINE)
        for i, (text, bold, color) in enumerate([
            (dim, True, ON_SURFACE), (before, False, ON_VARIANT), (after, True, ON_SURFACE)
        ]):
            tf = textbox(s, x + 16 + sum(col[:i]), ry + 13, col[i] - 20, 32)
            write(tf, text, 9, color, bold=bold, line=1.35, first=True)

    rect(s, x, 596, w, 74, fill=NAVY, radius=14)
    tf = textbox(s, x + 30, 620, 600, 30)
    write(tf, "VISIBILITY  →  CONTROL  →  IMPROVEMENT", 17, WHITE, bold=True,
          spacing=90, first=True)
    tf = textbox(s, x + 660, 616, w - 690, 40)
    write(tf, "Factory Vision MES: fondasi transformasi pabrik pintar menuju Industri 4.0.",
          9, NAVY_SOFT, align=PP_ALIGN.RIGHT, line=1.45, first=True)


for build in [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07,
              slide_08, slide_09, slide_10, slide_11, slide_12, slide_13, slide_14,
              slide_15, slide_16, slide_17, slide_18, slide_19, slide_20, slide_21,
              slide_22]:
    build()

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {OUT}")
